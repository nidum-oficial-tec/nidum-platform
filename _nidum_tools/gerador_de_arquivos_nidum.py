"""
title: Gerador de Arquivos Nidum
author: Nidum
version: 2.8.2
description: Gera PPTX, XLSX, DOCX, PDF, HTML e APRESENTACAO HTML navegavel no servidor com alto padrao de acabamento (UX/UI) e a identidade do brandbook Nidum: paleta, fonte Maxima Nouva embutida, logos, contraste correto, layouts variados, tabelas refinadas, rodapes e numeracao. Insere imagens anexadas pelo usuario. Devolve link de download nativo.
requirements: python-pptx, openpyxl, python-docx, reportlab
changelog:
  2.8.2:
    - A DOCSTRING ERA A CAUSA, O ALIAS E O SINTOMA. Telemetria de producao (4
      recusas identicas): o modelo mandou "corpo" e acertou TODO o resto que a
      docstring nomeia. A docstring dizia "corpo" 5x como nome do CONCEITO e
      "texto" 3x como nome do CAMPO - ele seguiu a palavra repetida. Reescrita
      para nomear o campo ("O corpo do slide se chama texto. Nao existe campo
      corpo") em vez de descrever o conceito. Mesma licao do chatnd.py:1420.
    - "corpo" acrescentado aos aliases de texto - o erro mais previsivel, e ficou
      de fora da primeira lista. Segue com telemetria e prazo.
    - Registrada como PROPRIEDADE DESEJADA (nao acidente) o fato de a mensagem de
      recusa ENSINAR: a recusa cai sempre no slide 3 e o modelo acerta do 4o em
      diante - uma recusa por geracao, nao uma por slide.
  2.8.1:
    - DEDUPE POR NOME no anexo nativo (medido em producao, 03/09/2026): o modelo
      chamou gerar_html duas vezes na mesma resposta - dois uuid, um nome so - e o
      dedupe por id nao pegava; o usuario via TRES cards (HTML, PPTX, HTML).
      Regerar o mesmo nome agora SUBSTITUI, mantendo a posicao original do card e
      usando a versao nova. Caso travado no teste_anexo_nativo.py.
  2.8.0:
    - ANEXO NATIVO: o arquivo deixa de passar pelo modelo. A tool emite
      {'type':'files'} pelo __event_emitter__ e o frontend faz message.files =
      data.files (Chat.svelte:523), renderizando o card de download. O retorno ao
      modelo vira uma frase SEM URL - nao ha o que corromper.
      POR QUE: duas correcoes de link fecharam duas portas e abriram outras (sumiu a
      barra inicial; depois o parentese do markdown vazou para dentro da URL). O
      problema nunca foi sintaxe, era a TRANSCRICAO. D22: o que precisa ser exato e
      imposto pelo codigo, nao pedido no texto.
    - ACUMULACAO POR MENSAGEM (requisito, nao borda). O frontend SUBSTITUI a lista a
      cada evento; "entregue em HTML e tambem em PPTX" - caso central da Fase D -
      faria o 2o arquivo apagar o 1o. Guarda por message_id e reemite a lista
      inteira. Teto de 64 mensagens em cache. Suite nova: teste_anexo_nativo.py.
    - _registrar_arquivo devolve COPIA da lista. Devolver a referencia viva fazia o
      payload do 1o evento ganhar o arquivo do 2o - invisivel hoje (o emitter
      serializa na hora), bomba-relogio se um dia nao serializar. Achado pelo teste.
    - CONTEXTVAR em vez de atributo de instancia: a Tools e reusada entre
      requisicoes, e guardar o emitter em self faria uma resposta emitir na mensagem
      de outra sob concorrencia.
    - FALLBACK DECLARADO: sem emitter (o pipe nunca passa um), volta o link - e o
      texto DIZ que e fallback, para aparecer na tela em vez de virar 404 mudo.
      A frase "Link para download:" e CONTRATO com o pipe (chatnd.py:5773 decide a
      oferta de outros formatos por ela) e ficou byte a byte igual.
    - WEBUI_URL vazia agora e REGISTRADA em log com a consequencia (link relativo),
      em vez de degradar calada. Configuracao vazia sem registro e como valve
      apontando para colecao apagada.
    - ACENTO EM ENUM DO MODELO. A validacao comparava tipo == "divisao" e o modelo
      mandou "divisao" COM ACENTO: nao casou, e o slide escapou da validacao que
      existia para pega-lo - o erro estava uma camada acima do que ele impedia.
      _fold (mesmo criterio do _f3_fold do pipe) normaliza 'tipo' e 'cor' UMA vez,
      em _normalizar_corpo_slides; os dois renderizadores passam a ver so a forma
      canonica. Varredura da classe inteira feita: sao estes dois campos, nao ha
      outros comparados por igualdade literal.
    - Os dois parametros novos entram no FIM da assinatura. Poe-los no meio quebrou
      quem chama por posicao (as suites de imagem passam 'imagens' assim) - pego
      pela suite antes de publicar.
  2.7.0:
    - CONSUMIDOR GENERATIVO NAS DUAS PONTAS. Os tres defeitos medidos no A.5 (03/09/2026)
      sao o MESMO deslocamento: a tool foi desenhada para o pipe, que montava a entrada
      por schema e repassava a saida verbatim. No laco agentico o MODELO escreve a
      entrada e transcreve a saida. O que precisa ser exato passa a ser imposto pelo
      codigo, nao pedido no texto.
    - LINK ABSOLUTO (defeito 1). _gravar resolve WEBUI_URL e devolve URL completa; o
      modelo reescrevia o caminho relativo sem a barra inicial e o navegador resolvia
      para /c/<chat_id>/api/... -> 404. Config vazia mantem o relativo (comportamento
      antigo). As 7 docstrings mandam copiar o link verbatim e dizem o que quebra.
      MEDIDO: pptx funcionou e html nao, na MESMA chamada - transcricao, nao codigo.
    - DOCSTRING DE gerar_pptx COM O SCHEMA QUE FUNCIONA (defeito 2, item principal).
      Portado do prompt do pipe (chatnd.py:1420), que e a versao provada: os 8 tipos
      de slide (eram 4), os campos de cada um, 'itens' e 'cor' (nao existiam), e a
      regra "o campo de corpo NUNCA pode vir vazio". EVIDENCIA: mesmo pedido, mesmo
      modelo, mesma base - o pipe fez 18 slides com corpo, o agente fez 10 mudos. A
      unica diferenca era a interface descrita.
    - ALIASES DE CORPO COM TELEMETRIA (defeito 2, rede). texto/bullets/itens aceitam
      nomes alternativos, e cada alias aceito e LOGADO com o nome que chegou. Tem
      PRAZO: uso zero por duas semanas -> saem. Alias sem telemetria vira esteira
      rolante. A tolerancia de FORMATO (string JSON) ja existia em _coerce/_lista.
    - VALIDACAO QUE GRITA (defeito 2). Slide de tipo que exige corpo (conteudo,
      destaque, divisao, numerada, cartoes) chegando vazio agora RECUSA com
      diagnostico e lista as chaves recebidas, em vez de gerar caixa muda em
      silencio. capa/secao/encerramento seguem validos sem corpo.
    - BLOCO-GUARDA DO CHROME (defeito 3). O CSS da marca entra antes de </head>, logo
      o CSS do modelo vinha depois e vencia por cascata (medido: "* {margin:0}"
      anulava o cabecalho, "img{width:100%}" esticava a logo). Um bloco com !important
      nas propriedades ESTRUTURAIS do chrome entra por ULTIMO, antes de </body>.
      ESCOPO DELIBERADO: conserta o CHROME, nao o CORPO - tema e cores do corpo sao
      escritos pelo modelo, e quem resolve isso e a skill de HTML.
    - @IMPORT EXTERNO REMOVIDO (defeito 3). Referencia remota de estilo escrita pelo
      modelo quebrava o offline que o base64 das fontes garante. Removida com log.
      NAO ha rejeicao de hex fora da paleta: exigiria parsear CSS e daria falso
      positivo em uso legitimo (branco de texto, cor de serie em grafico).
    - MaximaNouva-ExtraBold registrada. O arquivo esta no build desde 13/07/2026 e
      nunca era usada: faltava no @font-face do HTML, no do deck e no reportlab.
  2.6.0:
    - gerar_codigo (MODO PRESERVACAO; par do pipe chatnd 1.51.0). Grava codigo-fonte
      VERBATIM, sem injetar marca nem editor - para EDITAR o app do usuario (HTML com
      <script>/handlers) preservando o comportamento. Ao contrario do gerar_html, NAO
      chama _injetar_marca_html (repintaria o CSS do app) nem _injetar_editor (o
      contenteditable no body briga com os campos do formulario; o app ja tem os proprios
      controles). O que entra e o que sai, byte a byte.
    - NAO passa por _coerce: ele faz json.loads, e um .json literal viraria um dict
      reserializado (aspas trocadas) - o oposto de verbatim. Content-type por extensao
      (familia texto: html/htm/css/js/json/xml/md/txt/csv). Ext fora da familia -> .txt,
      com log (binario nao vem para ca - exige parser, fatia futura).
  2.5.0:
    - IMAGEM ANEXADA PELO USUARIO (capacidade NOVA; par do pipe chatnd 1.43.0 - republicar
      os dois juntos). Os metodos gerar_pptx, gerar_docx, gerar_pdf, gerar_html e
      gerar_apresentacao_html ganharam 'imagens: list = None' DEPOIS dos parametros
      existentes (chamadas atuais nao quebram). gerar_xlsx NAO recebe: imagem em planilha
      esta fora de escopo. Os bytes chegam por PARAMETRO, vindos do pipe - nunca por modelo.
    - COMO A POSICAO E DECIDIDA: o modelo poe um marcador ("IMAGEM_1") no campo 'imagem' do
      slide/secao; _imagem_do_item resolve pelo NUMERO (tolerante a "imagem 1", "1"). Em
      tipo=html nao ha campo, entao o marcador escrito no documento e substituido pela
      imagem (_inserir_imagens_html trata tanto <img src="IMAGEM_1"> quanto o marcador
      solto, e APAGA marcador orfao - senao "IMAGEM_1" sairia como texto cru no arquivo).
    - LAYOUT. PPTX e deck: a imagem ganha o PROPRIO slide logo apos o slide que a
      posicionou - os layouts (capa, cartoes, numerada, destaque) sao composicoes fechadas
      e encaixar foto dentro deles colidiria com texto; em slide proprio ela aparece
      grande, centralizada, com margem, com o titulo como legenda e a logo da marca.
      DOCX/PDF: inline, centralizada, no fim da secao. PROPORCAO E INVIOLAVEL: _encaixar
      calcula UM fator pelo lado limitante e aplica nos dois eixos (nunca estica nem
      achata); reduz o quanto precisar e amplia no maximo 2x (_MAX_AMPLIACAO), para nao
      deixar uma imagem pequena visivelmente pixelada.
    - DEGRADACAO SEGURA (filosofia do RAG opcional): imagem problematica NUNCA derruba a
      geracao. Formato detectado por MAGIC BYTES (jpeg/png/gif/webp), nao pelo mime
      declarado - que mente. Bytes que nao sao imagem, formato que a lib do formato de
      saida nao aceita (ex.: webp no PDF sem PIL com webp) ou falha na insercao: loga e o
      arquivo sai SEM aquela imagem. Item invalido nao DESLOCA os marcadores dos demais.
    - QUALIDADE: a imagem entra em RESOLUCAO ORIGINAL (bytes embutidos sem reamostrar nem
      recomprimir - PPTX/DOCX byte-identicos; PDF faz passthrough do JPEG/DCTDecode). A
      escala e do renderizador (PowerPoint/Word/leitor), em alta qualidade.
    - TODO (trade-off consciente, NAO mexer agora): como a imagem entra em resolucao
      original, um arquivo com varias fotos de celular pode ficar pesado (20 MB+) e
      atrapalhar e-mail/SharePoint. Fica a favor da qualidade por ora; encarar so quando
      incomodar (ex.: reamostrar acima de um teto de lado maior, com filtro LANCZOS).
  2.4.0:
    - APROVACAO DE MARCA (auditoria da 2.3.0). Fatia A (contraste/cores): a regra "pedra e
      SUPORTE, nunca texto sobre areia" (2.18:1) tinha sido consertada na 2.3.0 so no HTML/
      deck; estendida agora aos 3 formatos que faltavam - antetitulos do PPTX e assinaturas
      de rodape do PDF e DOCX: pedra 9D9890 -> escuro 1F1E1B. Bordas de tabela em cor nao-
      oficial DEDAD0 -> pedra 9D9890 (aqui pedra e uso CERTO: filete) no HTML/XLSX/PDF.
      Tons fora da paleta -> escuro: texto de citacao do HTML (5b574f) e palco do deck
      (15140f, decisao do dono: escuro oficial, nao um tom mais escuro que a paleta).
    - Fatia B (varredura de azul): procurado qualquer azul cravado a mao fora do ceu oficial
      4F7187 em todos os formatos - NENHUM encontrado. Os 7 azuis do arquivo sao todos
      4F7187. Sem mudanca de codigo; a varredura e o resultado.
    - Fatia C (assinatura do rodape): separador unificado em PONTO. HTML usava "&middot;",
      DOCX/PDF usavam " - "; agora os tres sao "nidum. fazer da casa um ninho." Wordmark
      'nidum' minusculo mantido no rodape; a frase e literal do Documento Fundador.
    - Fatia D (logo na abertura E no encerramento). Regra de marca: logo no topo e no fim,
      UMA vez cada (nao em toda pagina). Estado por formato: PPTX ja tinha (capa +
      encerramento); deck ja tinha (cover + logo do ultimo slide); PDF ja tinha no topo, so
      FALTAVA o fim (adicionado, centralizado). DOCX NAO tinha nenhum - adicionados topo e
      fim (centralizados, terracota sobre branco). HTML so tinha o rodape (fim) - adicionado
      o header (topo); e o rodape deixou de sumir no @media print, para a logo de fim
      sobreviver a impressao. XLSX nao recebe (grade, sem capa/fim). Logo terracota em fundo
      areia/branco (a logica de cor por fundo ja existente foi mantida).
  2.3.0:
    - EDITOR HTML embutido: barra fixa "Editar"/"Salvar HTML" no HTML e no deck gerados.
      Editar liga contenteditable no corpo (a barra fica de fora); Salvar serializa a pagina
      INTEIRA com as edicoes (document.documentElement.outerHTML, com <!DOCTYPE html>) e
      baixa via Blob - offline, sem servidor, e o arquivo baixado continua editavel. Ctrl/
      Cmd+S salva (previne o "salvar pagina" do navegador); some em @media print; injecao
      IDEMPOTENTE (marcador NIDUM_EDITOR - reabrir o arquivo salvo nao duplica a barra). No
      deck, o keydown das setas/espaco ignora quando a edicao esta ligada (window.__ndEdit),
      senao digitar trocava de slide. Valve EDITOR_HTML (default True) desliga. O nome
      sugerido no download e o _nome_padrao (calculado antes da montagem).
    - BUG gerar_html: o titulo era interpolado no <title> SEM escapar - '<' ou '&' quebravam
      a tag. Agora passa por _esc (o deck ja escapava).
    - NOMENCLATURA OFICIAL (Tec_Office_e_Governanca_de_Dados_v2, slide 7): arquivos agora
      saem como ECOSSISTEMA_TEMA_DD-MM-AAAA_vN.ext (ex.: MKT_Campanha_01-06-2026_v1.pptx),
      no lugar de titulo.replace(" ","_")+.ext (sem ecossistema, sem data, sem versao, com
      acento e barra crus). Helper _nome_padrao: fold de acento p/ ASCII, TEMA em CamelCase
      (<=60 chars), DATA no fuso de Brasilia (UTC-3 - o servidor roda em UTC; sem isto,
      arquivo gerado depois das 21h ganhava a data do dia seguinte), ECOSSISTEMA validado
      contra lista fechada (FONTE/REG/MKT/PROD/OPS/FIN/JUR/ACA/TEC/SUS/CC/CT/CE). Sigla
      vazia/invalida cai na valve ECOSSISTEMA_PADRAO (default TEC) e loga - o nome NUNCA
      derruba a geracao. Os 6 metodos ganharam ecossistema="" e versao=1 DEPOIS de __user__
      (nao quebra chamadas posicionais). Par: pipe chatnd 1.41.0 emite 'ecossistema' no
      GERADOR - republicar os dois juntos.
    - PALETA CORRIGIDA para o brandbook oficial (MKT_BrandbookNidum V1). Tres das seis
      cores estavam quase certas mas erradas, algumas cravadas a mao no CSS ignorando as
      constantes: musgo 647260 -> 515E52; pedra 8A8880 -> 9D9890; areia EAE6DC -> E5E0D5.
      Trocado em TODO lugar (constantes E hexes soltos, incl. rgba decimais). As outras
      tres ja estavam certas (terracota 9A4A2E, ceu 4F7187, escuro 1F1E1B).
    - CONTRASTE (efeito colateral da paleta): pedra 9D9890 sobre areia E5E0D5 da 2.18:1
      (reprova). Pedra e cor de SUPORTE no brandbook: rodape e antetitulo do deck, que a
      usavam como TEXTO, passaram a escuro suavizado rgba(31,30,27,.62); pedra fica so em
      filetes/bordas. E com o fundo virando areia, as superficies de alternancia do HTML
      (blockquote/code/pre) ficariam areia-sobre-areia (sumindo) - passaram a BRANCO. A
      zebra do XLSX/PDF NAO muda: assenta sobre branco (ROWBACKGROUNDS=[branco, cremealt]),
      onde areia contrasta - por isso CREME_ALT continua E5E0D5.
  2.2.0:
    - GERACAO NAO-BLOQUEANTE: a montagem dos arquivos (pptx/xlsx/docx/pdf/html),
      que e trabalho sincrono pesado de CPU, agora roda em thread separada
      (asyncio.to_thread). Antes, um PDF grande travava o event loop inteiro do
      Open WebUI e congelava TODOS os usuarios ate terminar. O upload ao Storage
      tambem foi para thread (importante quando o Storage virar rede, ex.: R2).
    - UPLOAD TOLERANTE A R2: o Cloudflare R2 nao implementa o header
      x-amz-tagging no PutObject e derruba uploads com tags (NotImplemented).
      O upload agora tenta com tags e, se falhar, repete SEM tags - o arquivo
      salva de qualquer forma. Prepara a migracao STORAGE_PROVIDER=s3 -> R2.
    - TRACEBACK SO NO LOG: erros internos nao despejam mais o traceback
      completo no chat do usuario (vazava caminhos e estrutura interna).
      Agora: log completo no servidor (logger "gerador_nidum") + mensagem
      limpa com CODIGO DE ERRO de 8 caracteres para correlacionar no log.
    - LIMPEZA: removidas as valves OPENWEBUI_BASE_URL e OPENWEBUI_API_KEY,
      que nunca eram usadas (campo de secret morto e passivo de seguranca).
"""

# ATENCAO: este arquivo usa apenas caracteres ASCII no codigo, de proposito.
# Nao insira bullets unicode, travessoes ou emojis aqui.
# Os textos com acento vem do usuario em tempo de execucao, o que nao e problema.

import asyncio
import contextvars
import io
import os
import json
import logging
import uuid
import inspect
import re
import datetime
import unicodedata

from pydantic import BaseModel

log = logging.getLogger("gerador_nidum")


# ----------------------------------------------------------------------------
# Nomenclatura oficial de arquivos (Tec_Office_e_Governanca_de_Dados_v2, slide 7):
#   ECOSSISTEMA_TEMA_DD-MM-AAAA_vN.ext   (ex.: MKT_CampanhaLancamento_01-06-2026_v1.pptx)
# NUNCA falha por causa do nome (o arquivo TEM que sair): ecossistema vazio/invalido cai
# no padrao (valve ECOSSISTEMA_PADRAO) e loga; acento e reduzido a ASCII.
# ----------------------------------------------------------------------------
_ECOSSISTEMAS = ("FONTE", "REG", "MKT", "PROD", "OPS", "FIN", "JUR", "ACA",
                 "TEC", "SUS", "CC", "CT", "CE")


def _fold_ascii(txt):
    # Acento quebra download em alguns clientes: reduz a ASCII. Preserva a CAIXA (o nome
    # usa CamelCase), ao contrario do _normalizar_ascii do pipe, que abaixa tudo.
    return (unicodedata.normalize("NFKD", txt or "")
            .encode("ascii", "ignore").decode("ascii"))


def _tema_camel(titulo, limite=60):
    # TEMA em CamelCase, so [A-Za-z0-9], cortado a ~limite chars.
    palavras = re.findall(r"[A-Za-z0-9]+", _fold_ascii(titulo or ""))
    camel = "".join(p[:1].upper() + p[1:] for p in palavras) or "Documento"
    return camel[:limite]


def _data_brasilia():
    # DD-MM-AAAA no fuso de Brasilia (UTC-3, fixo - o Brasil nao tem mais horario de
    # verao). O servidor roda em UTC; sem isto, um arquivo gerado depois das 21h (UTC)
    # ganharia a data do dia SEGUINTE.
    tz = datetime.timezone(datetime.timedelta(hours=-3))
    return datetime.datetime.now(tz).strftime("%d-%m-%Y")


def _nome_padrao(titulo, ecossistema, extensao, versao=1, padrao="TEC"):
    # Monta ECOSSISTEMA_TEMA_DATA_vN.ext. Robusto a lixo: sigla invalida -> padrao (+log);
    # versao nao-inteira -> 1; extensao com ou sem ponto. Nunca levanta por causa do nome.
    eco = _fold_ascii(str(ecossistema or "")).strip().upper()
    if eco not in _ECOSSISTEMAS:
        if eco:
            log.warning("gerador_nidum: ecossistema %r invalido - usando padrao %r",
                        ecossistema, padrao)
        eco = _fold_ascii(str(padrao or "TEC")).strip().upper()
        if eco not in _ECOSSISTEMAS:
            eco = "TEC"
    try:
        v = int(versao)
    except Exception:
        v = 1
    if v < 1:
        v = 1
    ext = str(extensao or "").lstrip(".").lower() or "bin"
    return "%s_%s_%s_v%d.%s" % (eco, _tema_camel(titulo), _data_brasilia(), v, ext)

# ----------------------------------------------------------------------------
# Identidade visual da Nidum (brandbook MKT) - cores em hex, sem o '#'
# ----------------------------------------------------------------------------
# Paleta oficial (MKT_BrandbookNidum V1). Nomes do brandbook entre parenteses.
NIDUM_TERRACOTA = "9A4A2E"   # terracota - assinatura / destaque
NIDUM_VERDE = "515E52"       # musgo - titulos e blocos
NIDUM_AZUL = "4F7187"        # ceu - blocos
NIDUM_CINZA = "9D9890"       # pedra - SUPORTE (filetes/bordas); NAO usar como texto
                             #         sobre areia (2.18:1 reprova) - ver rodape/ante
NIDUM_PRETO = "1F1E1B"       # escuro - cor de TEXTO (inclusive antetitulos)
NIDUM_CREME = "E5E0D5"       # areia - fundo principal
# CREME_ALT fica E5E0D5 de proposito: e a zebra de tabela do XLSX/PDF, que assenta sobre
# BRANCO (ROWBACKGROUNDS=[branco, cremealt]); ali areia contrasta. So no HTML - onde o
# fundo virou areia - a alternancia (blockquote/code/pre) passou a BRANCO. Nao unificar.
NIDUM_CREME_ALT = "E5E0D5"   # areia (zebra sobre branco no xlsx/pdf)
NIDUM_BRANCO = "FFFFFF"
# Tipografia da marca
NIDUM_FONT = "Maxima Nouva"  # titulos, subtitulos e corpo
NIDUM_FONT_LOGO = "Ibrand"   # exclusiva do logotipo
# Compatibilidade com versoes anteriores do gerador (nomes antigos -> paleta nova)
NIDUM_SALVIA = NIDUM_AZUL
NIDUM_OCRE = NIDUM_TERRACOTA
NIDUM_INK = NIDUM_PRETO
NIDUM_PAPER = NIDUM_CREME


def _hex_to_rgb(h):
    h = h.lstrip("#")
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _brand_dir():
    # Resolve a pasta de assets de marca (fontes + logos), embutida na imagem.
    candidates = []
    try:
        import open_webui

        candidates.append(
            os.path.join(os.path.dirname(open_webui.__file__), "static", "brand")
        )
    except Exception:
        pass
    candidates.append("/app/backend/open_webui/static/brand")
    for d in candidates:
        if d and os.path.isdir(d):
            return d
    return None


def _logo_path(cor):
    # cor: areia | azul | escuro | terracota | verde
    d = _brand_dir()
    if not d:
        return None
    p = os.path.join(d, "logos", "nidum-" + str(cor) + ".png")
    return p if os.path.isfile(p) else None


# --------------------------------------------------------- imagens enviadas pelo usuario
# Imagem que o USUARIO anexou no chat (2.5.0). Os bytes chegam do pipe por PARAMETRO
# (imagens=), nunca por modelo: uma foto em base64 vira um texto enorme no prompt (estoura
# contexto e provoca 429). O modelo so ve MARCADORES (IMAGEM_1, IMAGEM_2...) e diz ONDE
# cada uma entra, pelo campo 'imagem' do slide/secao. Custo de token: zero.
#
# FILOSOFIA (igual a do RAG opcional): imagem problematica NUNCA derruba a geracao. Se nao
# decodifica, se o formato nao e suportado pela lib do formato de saida, ou se a insercao
# falha - loga e segue SEM a imagem. O arquivo sai.
_MAX_AMPLIACAO = 2.0   # teto de ampliacao: acima disso, imagem pequena fica pixelada

_ASSINATURAS_IMG = (
    (b"\xff\xd8\xff", "jpeg"),
    (b"\x89PNG\r\n\x1a\n", "png"),
    (b"GIF87a", "gif"),
    (b"GIF89a", "gif"),
)


def _formato_imagem(b):
    # Detecta pelo CONTEUDO (magic bytes), nao pelo mime declarado - o mime do anexo
    # mente com frequencia. Devolve "" quando nao reconhece (o chamador degrada).
    if not b or len(b) < 12:
        return ""
    for assinatura, fmt in _ASSINATURAS_IMG:
        if b.startswith(assinatura):
            return fmt
    if b[0:4] == b"RIFF" and b[8:12] == b"WEBP":
        return "webp"
    return ""


def _decodificar_imagem(item):
    # Aceita bytes crus, data-URL base64 ("data:image/jpeg;base64,..."), base64 puro ou
    # dict {"url"/"dados"/"bytes"} - o formato que _extrair_imagens_anexo (pipe) devolve
    # e a data-URL. Devolve (bytes, formato); (None, "") quando nao da. NUNCA levanta.
    import base64

    try:
        if isinstance(item, dict):
            item = item.get("url") or item.get("dados") or item.get("bytes") or ""
        if isinstance(item, (bytes, bytearray)):
            b = bytes(item)
        elif isinstance(item, str):
            s = item.strip()
            if not s:
                return (None, "")
            if s.startswith("data:"):
                corte = s.find(",")
                if corte < 0:
                    return (None, "")
                s = s[corte + 1:]
            b = base64.b64decode(s)
        else:
            return (None, "")
        fmt = _formato_imagem(b)
        if not fmt:
            log.warning(
                "gerador_nidum: anexo descartado - formato de imagem nao reconhecido"
            )
            return (None, "")
        return (b, fmt)
    except Exception:
        log.exception("gerador_nidum: falha ao decodificar imagem anexada")
        return (None, "")


def _normalizar_imagens(imagens):
    # Lista vinda do pipe -> [{"marcador","bytes","formato"}]. A POSICAO define o
    # marcador (a 1a e IMAGEM_1), exatamente como o pipe conta ao modelo. Um item
    # invalido e descartado mas NAO desloca os demais - senao IMAGEM_2 viraria a 3a.
    saida = []
    for i, it in enumerate(_lista(imagens)):
        b, fmt = _decodificar_imagem(it)
        if not b:
            continue
        saida.append({"marcador": "IMAGEM_" + str(i + 1), "bytes": b, "formato": fmt})
    return saida


def _imagem_do_item(imagens_norm, item):
    # Resolve o campo 'imagem' de um slide/secao. Tolerante ao que o modelo escreve:
    # "IMAGEM_2", "imagem 2", "2" - vale o NUMERO. Sem casar -> None (sem imagem).
    if not imagens_norm or not isinstance(item, dict):
        return None
    ref = item.get("imagem")
    if ref is None or isinstance(ref, (list, dict)):
        return None
    m = re.search(r"(\d+)", str(ref))
    if not m:
        return None
    alvo = "IMAGEM_" + m.group(1)
    for img in imagens_norm:
        if img["marcador"] == alvo:
            return img
    return None


def _encaixar(nat_w, nat_h, box_w, box_h):
    # Escala UNIFORME nos dois eixos para caber na caixa: a imagem nunca estica nem
    # achata. Calcula pelo lado LIMITANTE (o menor fator) e deixa o outro seguir.
    # Reduz o quanto precisar; amplia no maximo _MAX_AMPLIACAO.
    try:
        nat_w = float(nat_w)
        nat_h = float(nat_h)
        if nat_w <= 0 or nat_h <= 0:
            return (box_w, box_h)
        escala = min(float(box_w) / nat_w, float(box_h) / nat_h)
        escala = min(escala, _MAX_AMPLIACAO)
        return (nat_w * escala, nat_h * escala)
    except Exception:
        return (box_w, box_h)


def _img_data_uri(img):
    # <img src="data:..."> para os formatos HTML (deck e pagina). Autocontido.
    import base64

    try:
        b64 = base64.b64encode(img["bytes"]).decode("ascii")
        return "data:image/" + img["formato"] + ";base64," + b64
    except Exception:
        log.exception("gerador_nidum: falha ao embutir imagem no html")
        return ""


def _inserir_imagens_html(c, imagens_norm):
    # Em tipo=html o modelo escreve HTML livre, entao nao ha campo 'imagem': ele posiciona
    # o MARCADOR no documento e aqui ele vira a imagem de verdade. Dois jeitos que o modelo
    # naturalmente escreve, ambos tratados:
    #   1) <img src="IMAGEM_1">  -> troca so o src (preserva a tag e os atributos dele)
    #   2) IMAGEM_1 solto        -> vira um <figure> centralizado com margem
    # Marcador sem imagem correspondente e APAGADO: melhor nada do que "IMAGEM_1" cru
    # aparecendo no documento final (era exatamente o sintoma do placeholder de texto).
    if not c:
        return c
    for img in imagens_norm or []:
        uri = _img_data_uri(img)
        if not uri:
            continue
        marc = img["marcador"]
        c = re.sub(
            r"""(<img\b[^>]*\bsrc\s*=\s*)(['"])\s*""" + marc + r"""\s*\2""",
            lambda m: m.group(1) + m.group(2) + uri + m.group(2),
            c,
            flags=re.IGNORECASE,
        )
        figura = (
            '<figure style="margin:28px auto;text-align:center;max-width:100%">'
            '<img src="' + uri + '" alt="" style="max-width:100%;height:auto;'
            'display:block;margin:0 auto;border-radius:10px">'
            "</figure>"
        )
        c = re.sub(r"\b" + marc + r"\b", figura, c, flags=re.IGNORECASE)
    # Varre marcadores orfaos (o modelo citou IMAGEM_9 sem existir) - nao deixa lixo.
    c = re.sub(r"\bIMAGEM_\d+\b", "", c, flags=re.IGNORECASE)
    return c


def _font_path(fname):
    d = _brand_dir()
    if not d:
        return None
    p = os.path.join(d, "fonts", fname)
    return p if os.path.isfile(p) else None


def _font_face(fname, weight, style="normal", family="Maxima Nouva"):
    # Gera uma regra @font-face com a fonte embutida em base64 (HTML autocontido).
    import base64

    p = _font_path(fname)
    if not p:
        return ""
    try:
        b64 = base64.b64encode(open(p, "rb").read()).decode("ascii")
    except Exception:
        log.exception("gerador_nidum: falha ao ler fonte %s", fname)
        return ""
    return (
        "@font-face{font-family:'%s';font-weight:%s;font-style:%s;"
        "src:url(data:font/ttf;base64,%s) format('truetype');}"
        % (family, weight, style, b64)
    )


def _brand_css():
    # Folha de estilo da marca Nidum para HTML/PDF (fontes embutidas + paleta).
    faces = "".join(
        [
            _font_face("MaximaNouva-Regular.ttf", 400),
            _font_face("MaximaNouva-SemiBold.ttf", 600),
            _font_face("MaximaNouva-Bold.ttf", 700),
            _font_face("MaximaNouva-ExtraBold.ttf", 800),
            _font_face("MaximaNouva-Italic.ttf", 400, "italic"),
            _font_face("Ibrand.ttf", "100 900", "normal", "Ibrand"),
        ]
    )
    rules = (
    # PALETA: os SEIS nomes oficiais do brandbook sao os canonicos; os antigos
    # viram ALIAS apontando para eles (--creme -> var(--areia) e assim por diante).
    # Alias, e nao troca seca, porque estes CSS ja sairam em arquivos que estao no
    # SharePoint e em caixas de e-mail: um .html gerado em agosto continua abrindo
    # com --creme, e quebrar isso seria estragar entrega passada para arrumar nome.
    # PRAZO: revisar em 2026-12 - se nenhuma regra nova usar os antigos ate la, os
    # aliases saem e ficam so os seis. Quem revisar: procure "--creme" e "--azul"
    # no repo; zero ocorrencias fora deste bloco = pode remover.
        ":root{--areia:#E5E0D5;--pedra:#9D9890;--terracota:#9A4A2E;"
        "--ceu:#4F7187;--musgo:#515E52;--escuro:#1F1E1B;"
        "--creme:var(--areia);--cinza:var(--pedra);--azul:var(--ceu);"
        "--verde:var(--musgo);--preto:var(--escuro);--cremealt:var(--areia);}"
        "*{box-sizing:border-box}"
        "html{background:#E5E0D5}"
        "body{background:#E5E0D5;color:#1F1E1B;"
        "font-family:'Maxima Nouva',-apple-system,Segoe UI,Roboto,Arial,sans-serif;"
        "line-height:1.72;margin:0 auto;max-width:880px;padding:64px 44px 72px;"
        "font-size:18px;-webkit-font-smoothing:antialiased;}"
        "h1{font-family:'Ibrand','Maxima Nouva',sans-serif;font-size:2.7em;"
        "color:#515E52;font-weight:400;line-height:1.08;"
        "letter-spacing:-.005em;margin:0 0 .5em;}"
        "h2{font-family:'Ibrand','Maxima Nouva',sans-serif;font-size:1.8em;"
        "color:#515E52;font-weight:400;line-height:1.18;"
        "margin:1.9em 0 .5em;padding-bottom:.22em;"
        "border-bottom:2px solid rgba(154,74,46,.22);}"
        "h3{font-size:1.28em;color:#9A4A2E;font-weight:600;margin:1.5em 0 .4em;}"
        "h4{font-size:1em;color:#515E52;text-transform:uppercase;"
        "letter-spacing:.14em;margin:1.5em 0 .4em;}"
        "p{margin:0 0 1.1em;}"
        "a{color:#9A4A2E;text-decoration:none;"
        "border-bottom:1px solid rgba(154,74,46,.4);}"
        "strong,b{color:#9A4A2E;font-weight:700;}"
        "ul,ol{margin:0 0 1.1em 1.3em;} li{margin:.45em 0;}"
        "blockquote{margin:1.7em 0;padding:.7em 1.5em;border-left:4px solid #9A4A2E;"
        "background:#FFFFFF;border-radius:0 12px 12px 0;color:#1F1E1B;font-style:italic;}"
        "hr{border:none;border-top:2px solid rgba(157,152,144,.32);margin:2.4em 0;}"
        "img{max-width:100%;height:auto;border-radius:14px;"
        "box-shadow:0 12px 32px rgba(31,30,27,.15);}"
        "table{border-collapse:separate;border-spacing:0;width:100%;margin:1.7em 0;"
        "border-radius:12px;overflow:hidden;box-shadow:0 6px 20px rgba(31,30,27,.08);}"
        "th{background:#515E52;color:#E5E0D5;text-align:left;padding:12px 14px;"
        "font-weight:600;}"
        "td{padding:11px 14px;border-bottom:1px solid #9D9890;}"
        "tr:nth-child(even) td{background:rgba(255,255,255,.45);}"
        "tr:last-child td{border-bottom:none;}"
        "code{font-family:Consolas,Menlo,monospace;background:#FFFFFF;"
        "border-radius:6px;padding:.12em .4em;font-size:.92em;}"
        "pre{font-family:Consolas,Menlo,monospace;background:#FFFFFF;"
        "border-radius:10px;padding:1em 1.2em;overflow:auto;}"
        ".nidum-header{margin:0 0 40px;}"
        ".nidum-header img{height:30px;box-shadow:none;border-radius:0;margin:0;}"
        ".nidum-footer{margin-top:64px;padding-top:18px;"
        "border-top:2px solid rgba(157,152,144,.3);display:flex;align-items:center;"
        "gap:10px;color:rgba(31,30,27,.62);font-size:.9em;letter-spacing:.02em;}"
        ".nidum-footer img{height:20px;box-shadow:none;border-radius:0;margin:0;}"
        "@media(max-width:900px){body{padding:40px 22px 60px;font-size:17px;}}"
        "@media print{body{max-width:none;padding:0;background:#fff;}}"
    )
    return "<style>/*NIDUM_BRAND*/\n" + faces + rules + "</style>"


_RE_IMPORT_EXTERNO = re.compile(
    r"@import\s+url\(\s*['\"]?https?://[^)]*\)\s*;?", re.IGNORECASE
)
_RE_LINK_EXTERNO = re.compile(
    r"<link\b[^>]*rel\s*=\s*['\"]?stylesheet['\"]?[^>]*https?://[^>]*>", re.IGNORECASE
)


def _guarda_chrome_css():
    # BLOCO-GUARDA (03/09/2026). O CSS da marca e inserido ANTES de </head>, entao o
    # CSS do modelo vem DEPOIS e vence por ordem de cascata - medido: um "* {margin:0}"
    # do modelo anulava o espacamento do cabecalho e um "img{width:100%}" esticava a
    # logo. Este bloco entra por ULTIMO (antes de </body>) e usa !important nas poucas
    # propriedades ESTRUTURAIS do chrome.
    #
    # ESCOPO, e ele e deliberado: isto conserta o CHROME (cabecalho, rodape, logo).
    # NAO conserta o CORPO. Se o modelo escrever o corpo em tema escuro com acento
    # fora da paleta, ele continua assim - quem escreve o corpo e o modelo, e quem
    # resolve isso e a skill de HTML, nao a ferramenta.
    return (
        "<style>/*NIDUM_GUARDA*/\n"
        "header.nidum-header,footer.nidum-footer{"
        "background:transparent!important;padding:24px 0!important;"
        "margin:0!important;border:0!important;display:block!important;"
        "text-align:center!important;}"
        "header.nidum-header img,footer.nidum-footer img{"
        "height:40px!important;width:auto!important;max-width:none!important;"
        "display:inline-block!important;object-fit:contain!important;"
        "filter:none!important;opacity:1!important;}"
        "footer.nidum-footer span{color:#9D9890!important;font-size:13px!important;"
        "display:block!important;margin-top:8px!important;}"
        "</style>"
    )


def _injetar_marca_html(conteudo):
    # Insere a folha de marca no <head> e um rodape sobrio antes de </body>.
    if "NIDUM_BRAND" in conteudo:
        return conteudo
    # OFFLINE: o base64 das fontes garante que o HTML funcione sem rede. Um @import
    # ou <link> remoto escrito pelo modelo quebra essa garantia (medido: o modelo
    # importou Inter do Google Fonts). Remove e segue - a substituta local assume.
    n_ext = len(_RE_IMPORT_EXTERNO.findall(conteudo)) + len(
        _RE_LINK_EXTERNO.findall(conteudo)
    )
    if n_ext:
        conteudo = _RE_IMPORT_EXTERNO.sub("", conteudo)
        conteudo = _RE_LINK_EXTERNO.sub("", conteudo)
        log.info(
            "gerador_nidum: %d referencia(s) externa(s) de estilo removida(s) "
            "(o HTML precisa abrir offline)", n_ext,
        )
    css = _brand_css()
    low = conteudo.lower()
    idx = low.find("</head>")
    if idx != -1:
        conteudo = conteudo[:idx] + css + conteudo[idx:]
    else:
        bidx = low.find("<body")
        if bidx != -1:
            end = conteudo.find(">", bidx)
            conteudo = (
                conteudo[: end + 1] + css + conteudo[end + 1 :]
                if end != -1
                else css + conteudo
            )
        else:
            conteudo = css + conteudo
    logo = _logo_b64("terracota")
    # Logo de ABERTURA (topo) - regra 2.4.0. Pagina areia -> logo terracota. Inserido logo
    # apos <body>; o rodape (abaixo) traz a logo de ENCERRAMENTO. Idempotente pelo marcador
    # NIDUM_BRAND (checado no topo desta funcao).
    if logo:
        header = "<header class='nidum-header'><img src='" + logo + "'></header>"
        lb = conteudo.lower().find("<body")
        if lb != -1:
            gt = conteudo.find(">", lb)
            if gt != -1:
                conteudo = conteudo[:gt + 1] + header + conteudo[gt + 1:]
    footer = (
        "<footer class='nidum-footer'>"
        + (("<img src='" + logo + "'>") if logo else "")
        + "<span>nidum. fazer da casa um ninho.</span></footer>"
    )
    # O bloco-guarda vai DEPOIS do rodape e ANTES de </body>: por ordem de cascata
    # ele e a ultima palavra sobre o chrome, mesmo que o modelo tenha escrito o
    # proprio CSS depois do <head>.
    fecho = footer + _guarda_chrome_css()
    bidx = conteudo.lower().rfind("</body>")
    if bidx != -1:
        conteudo = conteudo[:bidx] + fecho + conteudo[bidx:]
    else:
        conteudo = conteudo + fecho
    return conteudo


def _esc(t):
    return (
        str(t if t is not None else "")
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )


# ----------------------------------------------------------------------------
# EDITOR HTML embutido (v2.3.0): barra fixa "Editar" / "Salvar HTML" injetada no HTML e no
# deck gerados. Editar liga contenteditable no corpo (a barra fica DE FORA, senao o usuario
# editaria os botoes); Salvar serializa a PAGINA INTEIRA (com as edicoes) e baixa via Blob
# - offline, sem servidor, duplo clique - e o arquivo baixado CONTINUA editavel (a barra vai
# junto). Idempotente pelo marcador NIDUM_EDITOR: o arquivo salvo ja tem a barra e nao
# duplica se for reinjetado. Some em @media print. Rotulos ASCII (Editar/Salvar HTML).
# ----------------------------------------------------------------------------
_EDITOR_MARCADOR = "NIDUM_EDITOR"


def _esc_js(t):
    # Escapa uma string para caber dentro de um literal JS entre aspas duplas. O "</"
    # vira "<\/" para um nome com '</script>' nao fechar a tag no meio do bloco.
    return (str(t if t is not None else "")
            .replace("\\", "\\\\").replace('"', '\\"')
            .replace("\n", " ").replace("\r", " ").replace("</", "<\\/"))


def _bloco_editor_html(nome_download):
    nm = _esc_js(nome_download or "documento.html")
    css = (
        "<style>/*" + _EDITOR_MARCADOR + "*/"
        "#ndbar{position:fixed;top:14px;right:14px;z-index:2147483647;display:flex;gap:8px;"
        "font-family:-apple-system,Segoe UI,Roboto,Arial,sans-serif}"
        "#ndbar button{border:none;border-radius:999px;padding:8px 15px;font-size:14px;"
        "cursor:pointer;box-shadow:0 4px 14px rgba(31,30,27,.22);background:#515E52;"
        "color:#E5E0D5;font-weight:600;line-height:1}"
        "#ndbar button.alt{background:#9A4A2E;color:#E5E0D5}"
        "body.ndediting{outline:2px dashed rgba(154,74,46,.55);outline-offset:-8px}"
        "@media print{#ndbar{display:none!important}}"
        "</style>"
    )
    barra = (
        "<div id=\"ndbar\" contenteditable=\"false\">"
        "<button id=\"ndedit\" type=\"button\">Editar</button>"
        "<button id=\"ndsave\" type=\"button\" class=\"alt\">Salvar HTML</button>"
        "</div>"
    )
    js = (
        "<script>/*" + _EDITOR_MARCADOR + "*/(function(){"
        "if(window.__ndInit){return;}window.__ndInit=true;"
        "var NM=\"" + nm + "\";window.__ndEdit=false;"
        "function bar(){return document.getElementById('ndbar');}"
        "function setEdit(on){window.__ndEdit=!!on;"
        "document.body.setAttribute('contenteditable',on?'true':'false');"
        "document.body.classList.toggle('ndediting',!!on);"
        "var bb=bar();if(bb){bb.setAttribute('contenteditable','false');}"
        "var e=document.getElementById('ndedit');if(e){e.textContent=on?'Editando':'Editar';}}"
        "function save(){var was=window.__ndEdit;if(was){setEdit(false);}"
        "var html='<!DOCTYPE html>\\n'+document.documentElement.outerHTML;"
        "var blob=new Blob([html],{type:'text/html;charset=utf-8'});"
        "var a=document.createElement('a');a.href=URL.createObjectURL(blob);a.download=NM;"
        "document.body.appendChild(a);a.click();"
        "setTimeout(function(){URL.revokeObjectURL(a.href);a.remove();},1500);"
        "if(was){setEdit(true);}}"
        "var eb=document.getElementById('ndedit');"
        "if(eb){eb.addEventListener('click',function(){setEdit(!window.__ndEdit);});}"
        "var sb=document.getElementById('ndsave');"
        "if(sb){sb.addEventListener('click',save);}"
        "document.addEventListener('keydown',function(ev){"
        "if((ev.ctrlKey||ev.metaKey)&&(ev.key==='s'||ev.key==='S')){ev.preventDefault();save();}"
        "},true);"
        "})();</script>"
    )
    return css + barra + js


def _injetar_editor(html, nome_download):
    # Insere a barra antes de </body>. Idempotente: se o marcador ja existe (arquivo salvo
    # sendo reinjetado), devolve intacto - nunca duplica a barra.
    if not html:
        return html or ""
    if _EDITOR_MARCADOR in html:
        return html
    bloco = _bloco_editor_html(nome_download)
    idx = html.lower().rfind("</body>")
    if idx == -1:
        return html + bloco
    return html[:idx] + bloco + html[idx:]


def _logo_b64(cor):
    import base64

    p = _logo_path(cor)
    if not p:
        return ""
    try:
        return "data:image/png;base64," + base64.b64encode(
            open(p, "rb").read()
        ).decode("ascii")
    except Exception:
        log.exception("gerador_nidum: falha ao ler logo %s", cor)
        return ""


DECK_CSS = (
    "*{box-sizing:border-box;margin:0;padding:0}"
    # Mesma paleta do bloco acima (nomes oficiais + alias). Os dois blocos sao
    # copias por desenho - CSS embutido em arquivo entregue nao pode depender de
    # import; se um mudar, o outro muda junto. O teste afirma que sao iguais.
    ":root{--areia:#E5E0D5;--pedra:#9D9890;--terracota:#9A4A2E;"
    "--ceu:#4F7187;--musgo:#515E52;--escuro:#1F1E1B;"
    "--creme:var(--areia);--cinza:var(--pedra);--azul:var(--ceu);"
    "--verde:var(--musgo);--preto:var(--escuro);--cremealt:var(--areia);"
    "--rn:26px;"
    "--sc:0 26px 70px rgba(31,30,27,.34)}"
    "html,body{height:100%;background:#1F1E1B;"
    "font-family:'Maxima Nouva',-apple-system,Segoe UI,Roboto,Arial,sans-serif}"
    ".deck{position:fixed;inset:0;display:flex;align-items:center;"
    "justify-content:center;padding:3vh 3vw}"
    ".slide{position:absolute;width:min(94vw,1180px);aspect-ratio:16/9;"
    "border-radius:var(--rn);overflow:hidden;box-shadow:var(--sc);opacity:0;"
    "transform:translateY(16px) scale(.985);transition:opacity .5s ease,"
    "transform .55s cubic-bezier(.2,.7,.2,1);pointer-events:none;display:flex;"
    "flex-direction:column;justify-content:center;padding:6% 7%;background:#E5E0D5;"
    "color:#1F1E1B}"
    ".slide.active{opacity:1;transform:none;pointer-events:auto}"
    ".slide .ante{font-size:.92rem;letter-spacing:.2em;text-transform:uppercase;"
    "color:rgba(31,30,27,.62);margin-bottom:1rem}"
    ".slide h1{font-family:'Ibrand','Maxima Nouva',sans-serif;"
    "font-size:clamp(2.1rem,5.4vw,3.9rem);color:var(--verde);"
    "font-weight:400;line-height:1.05}"
    ".slide h2{font-family:'Ibrand','Maxima Nouva',sans-serif;"
    "font-size:clamp(1.6rem,3.8vw,2.6rem);color:var(--verde);"
    "font-weight:400;line-height:1.1}"
    ".slide p{font-size:clamp(1rem,1.7vw,1.28rem);line-height:1.62;margin-top:1rem;"
    "max-width:60ch;color:var(--preto)}"
    ".slide ul{margin:1.1rem 0 0 1.2rem} .slide li{font-size:clamp(1rem,1.6vw,1.2rem);"
    "line-height:1.5;margin:.55rem 0}"
    ".slide .logo{height:30px;position:absolute;right:6%;bottom:6.4%;opacity:.92}"
    ".slide.cover{align-items:center;text-align:center}"
    ".slide.cover .logo-c{height:clamp(54px,8vw,108px);margin-bottom:1.4rem}"
    ".slide.cover .sub{color:var(--terracota);font-size:clamp(1.1rem,2vw,1.5rem);"
    "margin-top:.7rem}"
    ".slide.fill{justify-content:center;color:var(--creme)}"
    ".slide.fill .ante{color:rgba(229,224,213,.72)}"
    ".slide.fill h1,.slide.fill h2{color:var(--creme)} .slide.fill p{color:var(--creme)}"
    ".slide.split{flex-direction:row;padding:0}"
    ".split .left{flex:0 0 42%;display:flex;align-items:center;padding:6.5%}"
    ".split .left h2{color:var(--creme)}"
    ".split .right{flex:1;background:var(--creme);display:flex;flex-direction:column;"
    "justify-content:center;padding:6.5%}"
    ".num{display:flex;gap:1.1rem;align-items:baseline;margin:.85rem 0}"
    ".num .n{font-size:clamp(1.8rem,3vw,2.7rem);color:var(--terracota);"
    "font-weight:700;line-height:1;min-width:1.5em}"
    ".num .t{font-size:clamp(1rem,1.6vw,1.22rem);line-height:1.45}"
    ".num .t b{color:var(--verde)}"
    ".cards{display:grid;grid-template-columns:1fr 1fr;gap:1.1rem;margin-top:1.3rem}"
    ".card{border-radius:18px;padding:1.3rem 1.45rem;color:var(--creme)}"
    ".card h3{font-size:1.22rem;margin-bottom:.35rem} "
    ".card p{color:rgba(229,224,213,.92);margin-top:.1rem;font-size:1.02rem;max-width:none}"
    ".nav{position:fixed;bottom:20px;left:50%;transform:translateX(-50%);display:flex;"
    "align-items:center;gap:14px;background:rgba(31,30,27,.55);"
    "-webkit-backdrop-filter:blur(7px);backdrop-filter:blur(7px);padding:8px 16px;"
    "border-radius:999px;z-index:10}"
    ".nav button{border:none;background:transparent;color:#E5E0D5;font-size:1.3rem;"
    "cursor:pointer;width:34px;height:34px;border-radius:50%;line-height:1;"
    "transition:background .15s}"
    ".nav button:hover{background:rgba(229,224,213,.18)}"
    ".dots{display:flex;gap:7px}"
    ".dot{width:8px;height:8px;border-radius:50%;background:rgba(229,224,213,.42);"
    "cursor:pointer;transition:background .2s,width .2s}"
    ".dot.on{background:#E5E0D5;width:22px;border-radius:999px}"
    ".count{color:#E5E0D5;font-size:.84rem;min-width:52px;text-align:center;"
    "letter-spacing:.04em}"
    "@media print{.nav{display:none}.slide{position:relative;opacity:1!important;"
    "transform:none!important;margin:0 auto 24px;page-break-after:always}}"
)

DECK_JS = (
    "<script>var S=[].slice.call(document.querySelectorAll('.slide')),"
    "D=[].slice.call(document.querySelectorAll('.dot')),i=0;"
    "function go(n){i=Math.max(0,Math.min(S.length-1,n));"
    "S.forEach(function(s,k){s.classList.toggle('active',k===i);});"
    "D.forEach(function(d,k){d.classList.toggle('on',k===i);});"
    "var c=document.getElementById('cnt');if(c){c.textContent=(i+1)+' / '+S.length;}}"
    "document.addEventListener('keydown',function(e){"
    "if(window.__ndEdit){return;}"
    "if(e.key==='ArrowRight'||e.key==='PageDown'||e.key===' '){e.preventDefault();go(i+1);}"
    "if(e.key==='ArrowLeft'||e.key==='PageUp'){e.preventDefault();go(i-1);}});"
    "go(0);</script>"
)


def _slide_html(s, tipo, mapa, cores_secao, cores_cartao, sec, logo_t, logo_a):
    titulo = _esc(s.get("titulo"))
    sub = _esc(s.get("subtitulo"))
    texto = _esc(s.get("texto"))
    logo_creme = ("<img class='logo' src='" + logo_t + "'>") if logo_t else ""
    logo_color = ("<img class='logo' src='" + logo_a + "'>") if logo_a else ""

    def bl():
        items = _lista(s.get("bullets"))
        if not items:
            return ""
        return "<ul>" + "".join("<li>" + _esc(b) + "</li>" for b in items) + "</ul>"

    ante = ("<div class='ante'>" + sub + "</div>") if sub else ""
    h1 = ("<h1>" + titulo + "</h1>") if titulo else ""
    h2 = ("<h2>" + titulo + "</h2>") if titulo else ""
    par = ("<p>" + texto + "</p>") if texto else ""

    if tipo == "capa":
        lg = ("<img class='logo-c' src='" + logo_t + "'>") if logo_t else ""
        subc = ("<div class='sub'>" + sub + "</div>") if sub else ""
        return "<section class='slide cover'>" + lg + h1 + subc + "</section>"
    if tipo == "secao":
        c = mapa.get(str(s.get("cor", "")).lower(), cores_secao[sec[0] % len(cores_secao)])
        sec[0] += 1
        return ("<section class='slide fill' style='background:" + c + "'>"
                + ante + h1 + logo_color + "</section>")
    if tipo == "destaque":
        c = mapa.get(str(s.get("cor", "terracota")).lower(), "#9A4A2E")
        return ("<section class='slide fill' style='background:" + c + "'>"
                + ante + h1 + par + logo_color + "</section>")
    if tipo == "divisao":
        c = mapa.get(str(s.get("cor", "verde")).lower(), "#515E52")
        right = ante + par + bl()
        return ("<section class='slide split'><div class='left' style='background:"
                + c + "'>" + h2 + "</div><div class='right'>" + right + "</div></section>")
    if tipo == "numerada":
        itens = _lista_de_dicts(s.get("itens"))
        norm = []
        for it in itens:
            norm.append(
                (_esc(it.get("titulo")),
                 _esc(it.get("texto") or " ".join(_lista(it.get("bullets")))))
            )
        if not norm:
            norm = [("", _esc(b)) for b in _lista(s.get("bullets"))]
        rows = ""
        for k, (t_i, x_i) in enumerate(norm[:6]):
            inner = ("<b>" + t_i + "</b> " + x_i) if t_i else x_i
            rows += ("<div class='num'><div class='n'>" + str(k + 1)
                     + "</div><div class='t'>" + inner + "</div></div>")
        return "<section class='slide'>" + ante + h2 + rows + logo_creme + "</section>"
    if tipo == "cartoes":
        itens = _lista_de_dicts(s.get("itens"))
        if not itens:
            itens = [{"texto": b} for b in _lista(s.get("bullets"))]
        cards = ""
        for k, it in enumerate(itens[:4]):
            c = cores_cartao[k % len(cores_cartao)]
            ch = ("<h3>" + _esc(it.get("titulo")) + "</h3>") if it.get("titulo") else ""
            cp = ("<p>" + _esc(it.get("texto")) + "</p>") if it.get("texto") else ""
            cards += ("<div class='card' style='background:" + c + "'>" + ch + cp + "</div>")
        return ("<section class='slide'>" + ante + h2 + "<div class='cards'>"
                + cards + "</div>" + logo_creme + "</section>")
    if tipo == "encerramento":
        lg = ("<img class='logo-c' src='" + logo_t + "'>") if logo_t else ""
        subc = ("<div class='sub'>" + sub + "</div>") if sub else ""
        return "<section class='slide cover'>" + lg + h1 + subc + "</section>"
    return ("<section class='slide'>" + ante + h2 + par + bl()
            + logo_creme + "</section>")


def _get_user_id(__user__):
    if __user__ is None:
        return None
    if isinstance(__user__, dict):
        return __user__.get("id")
    return getattr(__user__, "id", None)


def _coerce(value):
    # O modelo as vezes envia listas/objetos como string JSON. Converte de volta.
    if isinstance(value, str):
        try:
            return json.loads(value)
        except Exception:
            return value
    return value


def _lista_de_dicts(value):
    # Garante uma lista de dicionarios, aceitando string JSON ou dict unico.
    value = _coerce(value)
    if isinstance(value, dict):
        value = [value]
    out = []
    for item in value or []:
        item = _coerce(item)
        if isinstance(item, dict):
            out.append(item)
    return out


def _lista(value):
    # Garante uma lista simples, aceitando string JSON.
    value = _coerce(value)
    if value is None:
        return []
    if isinstance(value, (list, tuple)):
        return list(value)
    return [value]


def _quebras(txt):
    # SEPARADOR DE PARAGRAFO -> QUEBRA NORMAL.
    #
    # Texto vindo de documento (Word, PowerPoint, PDF convertido) usa VT (\x0b)
    # onde o autor apertou Shift+Enter, e FF (\x0c) na quebra de pagina. O Word
    # chama os dois de "quebra"; o Python nao - str.split("\n") ignora ambos.
    #
    # O efeito era MUDO, que e o pior tipo: o bloco virava UMA linha so, o separador
    # sobrava dentro do texto como caractere de controle, e o PowerPoint desenhava um
    # retangulo vazio no meio da frase. Sem erro, sem log - um slide feio que ninguem
    # liga ao caractere que o causou.
    #
    # Normaliza UMA vez, na fronteira: daqui para dentro so existe "\n".
    # \r\n e \r entram junto por terem a mesma origem (arquivo que passou por
    # Windows ou por Mac antigo); LS/PS (\u2028/\u2029) pela mesma razao, vindos de
    # copiar-colar de PDF.
    t = str(txt or "")
    for bruto in ("\r\n", "\r", "\x0b", "\x0c", "\u2028", "\u2029"):
        t = t.replace(bruto, "\n")
    return t

def _texto_para_slide(txt):
    # Converte uma string "Titulo\n- bullet\n- bullet\ntexto" em slide dict.
    linhas = [l.strip() for l in _quebras(txt).split("\n") if l.strip()]
    titulo = ""
    bullets = []
    textos = []
    for idx, l in enumerate(linhas):
        if idx == 0 and l[:1] not in ("-", "*"):
            titulo = l
        elif l[:1] in ("-", "*"):
            bullets.append(l[1:].strip())
        else:
            textos.append(l)
    slide = {"tipo": "conteudo", "titulo": titulo}
    if bullets:
        slide["bullets"] = bullets
    if textos:
        slide["texto"] = " ".join(textos)
    return slide


def _texto_para_secao(txt):
    # Converte uma string "Heading\n- bullet\nparagrafo" em secao dict.
    linhas = [l.strip() for l in _quebras(txt).split("\n") if l.strip()]
    heading = ""
    paragrafos = []
    bullets = []
    for idx, l in enumerate(linhas):
        if idx == 0 and l[:1] not in ("-", "*"):
            heading = l
        elif l[:1] in ("-", "*"):
            bullets.append(l[1:].strip())
        else:
            paragrafos.append(l)
    sec = {"heading": heading}
    if paragrafos:
        sec["paragrafos"] = paragrafos
    if bullets:
        sec["bullets"] = bullets
    return sec


def _itens_loose(value, conversor_string):
    # Aceita lista de dicts OU lista de strings OU dict/string unico.
    # Strings sao convertidas em dict pelo conversor informado. Isso torna a
    # tool tolerante ao jeito que o modelo costuma enviar (lista de textos).
    value = _coerce(value)
    if isinstance(value, (dict, str)):
        value = [value]
    out = []
    for item in value or []:
        item = _coerce(item)
        if isinstance(item, dict):
            out.append(item)
        elif isinstance(item, str) and item.strip():
            out.append(conversor_string(item))
    return out


# ---- CORPO DO SLIDE: aliases com telemetria + validacao alta (03/09/2026) ----
# POR QUE EXISTE: no laco agentico o MODELO escreve a entrada da ferramenta, e a
# docstring e a unica interface que ele le. Medido em 03/09: o pipe (schema por
# extenso no prompt) gerou 18 slides com corpo cheio; o agente, mesmo pedido e
# mesmo modelo, 10 slides MUDOS. A docstring foi corrigida com o schema que
# funciona; estes aliases sao REDE, nao conserto.
#
# PRAZO: cada alias aceito e logado com o nome que chegou. Se o log mostrar uso
# ZERO por duas semanas, os aliases SAEM - alias sem telemetria vira esteira
# rolante, e o que sustenta a remocao e o dado, nao a lembranca de alguem.
#
# O QUE A TELEMETRIA JA ENSINOU (medido em producao, 03/09/2026): quatro recusas,
# todas com as mesmas chaves - 'corpo, destaque_rotulo, destaque_texto, rotulo,
# tipo, titulo'. O modelo ACERTOU tudo que a docstring nomeia e errou so o corpo.
# A causa nao era vocabulario dele: era a MINHA docstring, que dizia "corpo" cinco
# vezes como o nome do conceito e "'texto'" tres como o nome do campo. Ele seguiu a
# palavra repetida. A docstring foi corrigida para nomear o campo em vez de descrever
# o conceito ("O corpo do slide se chama 'texto'. Nao existe campo 'corpo'").
# 'corpo' entra como alias porque e o erro mais previsivel de todos - mas o conserto
# de verdade foi o nome, nao o alias. Alias trata sintoma; nome claro trata a causa.
_ALIAS_TEXTO = ("texto", "corpo", "conteudo", "body", "paragrafo", "paragrafos",
                "descricao")
_ALIAS_BULLETS = ("bullets", "pontos", "itens_texto", "topicos", "lista", "marcadores")
_ALIAS_ITENS = ("itens", "items", "cartoes", "etapas")
# Tipos que ficam MUDOS sem corpo. 'capa', 'secao' e 'encerramento' sao legitimos
# sem corpo - nao entram aqui.
_TIPOS_EXIGEM_CORPO = ("conteudo", "destaque", "divisao", "numerada", "cartoes")


def _campo_alias(slide, aliases, metodo="gerar_pptx"):
    # Devolve o valor do primeiro alias presente. Loga quando o nome usado NAO foi
    # o canonico (o primeiro da tupla) - e essa linha que decide se os aliases ficam.
    canonico = aliases[0]
    for nome in aliases:
        valor = slide.get(nome)
        if valor:
            if nome != canonico:
                log.info(
                    "gerador_nidum: %s ALIAS ACEITO '%s' -> '%s' (slide tipo=%s)",
                    metodo, nome, canonico, slide.get("tipo") or "?",
                )
            return valor
    return None


def _diag_slide_mudo(metodo, idx, slide):
    # Slide de tipo que EXIGE conteudo e chegou sem nenhum. Antes isto gerava uma
    # caixa vazia no arquivo, em silencio - o defeito medido em 03/09. Agora grita, e
    # diz QUAIS chaves chegaram: sem isso, o proximo diagnostico volta a ser
    # adivinhacao.
    #
    # PROPRIEDADE DESEJADA, NAO ACIDENTE: a mensagem ENSINA. Medido em producao -
    # a recusa acontece sempre no slide 3 (o primeiro de conteudo, depois de capa e
    # sumario), o modelo le o erro, corrige, e acerta do quarto em diante. Uma
    # recusa por geracao, nao uma por slide. Por isso a mensagem nomeia os campos
    # certos E lista os que chegaram: recusar sem ensinar geraria o mesmo erro em
    # laco ate estourar o teto de iteracoes. Toda validacao nova desta ferramenta
    # deve manter isso - dizer o que falta E o que veio no lugar.
    chaves = sorted([str(k) for k in (slide or {}).keys()])
    log.warning(
        "gerador_nidum: %s slide %d tipo=%s SEM CORPO. Chaves recebidas: %s",
        metodo, idx, slide.get("tipo") or "?", ", ".join(chaves) or "(nenhuma)",
    )
    return (
        "DIAGNOSTICO " + metodo + ": o slide " + str(idx) + " (tipo '"
        + str(slide.get("tipo") or "?") + "') chegou sem nenhum campo de corpo, "
        "e sairia mudo no arquivo. Esse tipo exige 'texto', 'bullets' ou 'itens'. "
        "As chaves que chegaram foram: " + (", ".join(chaves) or "(nenhuma)") + ". "
        "Reenvie o slide com o corpo preenchido."
    )


# ---- ANEXO NATIVO (03/09/2026): o arquivo aparece SEM passar pelo modelo ----
# POR QUE: duas correcoes de link fecharam duas portas e abriram outras - primeiro a
# barra inicial sumia (/c/api/...), depois o parentese do markdown vazava para dentro
# da URL. O problema nunca foi sintaxe: e a TRANSCRICAO. Instrucao em docstring e
# pedido; o D22 manda impor pelo codigo.
#
# COMO: o evento {'type':'files'} faz o frontend executar message.files = data.files
# (Chat.svelte:523) e renderizar o card de download nativo. O modelo nao ve endereco
# nenhum - recebe so uma frase de confirmacao.
#
# ACUMULACAO E REQUISITO, NAO DETALHE: o frontend SUBSTITUI a lista a cada evento.
# Num pedido como "entregue em HTML e tambem em PPTX" - caso central da Fase D - o
# segundo evento apagaria o primeiro da mensagem. Por isso guardamos por message_id e
# reemitimos a lista INTEIRA a cada arquivo novo. Trocar um bug visivel (404) por um
# invisivel (arquivo que some) seria pior que o estado atual.
_ARQUIVOS_POR_MENSAGEM = {}
_MAX_MENSAGENS_EM_CACHE = 64

# CONTEXTVAR, e nao atributo de instancia: a Tools e reusada entre requisicoes, e
# guardar o emitter em self faria uma resposta emitir na mensagem de outra sob
# concorrencia. ContextVar e por tarefa asyncio - isolada por definicao. Evita
# tambem passar dois parametros por 14 assinaturas internas.
_CTX_EMITTER = contextvars.ContextVar("nidum_emitter", default=None)
_CTX_MESSAGE_ID = contextvars.ContextVar("nidum_message_id", default=None)


def _ctx_entrega(emitter, message_id):
    # Chamado no inicio de cada metodo publico, com o que o Open WebUI injetou.
    _CTX_EMITTER.set(emitter)
    _CTX_MESSAGE_ID.set(message_id)


def _registrar_arquivo(message_id, item):
    # Acumula e devolve a lista completa da mensagem. Sem message_id (chamada fora de
    # um chat) nao ha o que acumular: devolve so o item.
    if not message_id:
        return [item]
    lista = _ARQUIVOS_POR_MENSAGEM.setdefault(message_id, [])
    # DEDUPE POR NOME, nao por id (medido em producao 03/09/2026): quando o modelo
    # chama gerar_html duas vezes na mesma resposta, saem DOIS arquivos com uuid
    # diferente e o MESMO nome - dedupe por id nao pega, e o usuario ve tres cards
    # (HTML, PPTX, HTML). Regenerar o mesmo nome e SUBSTITUIR, nao acrescentar.
    # Mantem a POSICAO original: o card nao pula de lugar quando e regerado.
    for idx, existente in enumerate(lista):
        if existente.get("name") == item.get("name"):
            if existente.get("id") != item.get("id"):
                log.info(
                    "gerador_nidum: %s regerado na mesma mensagem - card substituido "
                    "(id %s -> %s)", item.get("name"), existente.get("id"),
                    item.get("id"),
                )
            lista[idx] = item
            break
    else:
        lista.append(item)
    # Teto de memoria: o dict e de processo e viveria para sempre sem isto.
    if len(_ARQUIVOS_POR_MENSAGEM) > _MAX_MENSAGENS_EM_CACHE:
        for k in list(_ARQUIVOS_POR_MENSAGEM)[:-_MAX_MENSAGENS_EM_CACHE]:
            _ARQUIVOS_POR_MENSAGEM.pop(k, None)
    # COPIA, nao a lista viva: o evento ja emitido nao pode mudar depois. Devolver a
    # referencia fazia o payload do 1o evento ganhar o arquivo do 2o - invisivel em
    # producao (o emitter serializa na hora) e uma bomba-relogio se um dia nao
    # serializar. Achado pelo teste_anexo_nativo.
    return list(lista)


async def _emitir_arquivo(emitter, message_id, file_id, filename, content_type, url):
    # Devolve True se o anexo nativo foi entregue. False -> o chamador avisa que
    # esta caindo no fallback (nunca falhar em silencio: e o que produziu o 404).
    if emitter is None:
        return False
    item = {
        "type": "file",
        "id": file_id,
        "name": filename,
        "url": url,
        "content_type": content_type,
    }
    try:
        resultado = emitter(
            {"type": "files", "data": {"files": _registrar_arquivo(message_id, item)}}
        )
        if inspect.isawaitable(resultado):
            await resultado
        return True
    except Exception:
        log.exception("gerador_nidum: falha ao emitir o anexo nativo de %s", filename)
        return False


def _fold(s):
    # Dobra acento e caixa. MESMO criterio do _f3_fold do pipe - reusado de proposito:
    # os dois recebem texto escrito por modelo, e ter dois criterios de comparacao para
    # o mesmo dado e como ter dois relogios.
    s = unicodedata.normalize("NFD", str(s or "").strip())
    return "".join(c for c in s if not unicodedata.combining(c)).lower()


# Campos de ENUM que chegam do modelo e sao comparados por igualdade literal no
# codigo. Todo campo aqui e dobrado UMA vez, em _normalizar_corpo_slides, e os
# renderizadores passam a ver so a forma canonica.
#
# POR QUE ISTO EXISTE (03/09/2026): a validacao de slide mudo comparava
# tipo == "divisao" e o modelo mandou "divisao" COM ACENTO - nao casou, e o slide
# escapou da validacao que existia para pega-lo. O erro estava uma camada acima do
# que ele deveria impedir: comparacao literal sobre texto de modelo. Varredura feita
# na mesma sessao: 'tipo' e comparado em DOIS renderizadores (_slide_html e _pptx) e
# 'cor' resolvia por .lower() sem dobrar acento. Sao estes dois - nao ha outros.
_ENUMS_DO_MODELO = ("tipo", "cor")


def _normalizar_corpo_slides(slides, metodo="gerar_pptx"):
    # Resolve os aliases para os nomes canonicos e valida o corpo. Devolve
    # (slides, erro): erro != None significa PARE e devolva o erro ao chamador -
    # e melhor recusar com diagnostico do que entregar um arquivo mudo.
    for idx, s in enumerate(slides or [], start=1):
        if not isinstance(s, dict):
            continue
        # Enums dobrados ANTES de qualquer comparacao - os renderizadores comparam
        # por igualdade literal e passam a ver so a forma canonica.
        for campo in _ENUMS_DO_MODELO:
            if s.get(campo):
                bruto = str(s[campo])
                s[campo] = _fold(bruto)
                if s[campo] != bruto.strip().lower():
                    log.info(
                        "gerador_nidum: %s campo '%s' normalizado %r -> %r",
                        metodo, campo, bruto, s[campo],
                    )
        for aliases in (_ALIAS_TEXTO, _ALIAS_BULLETS, _ALIAS_ITENS):
            valor = _campo_alias(s, aliases, metodo)
            if valor and not s.get(aliases[0]):
                s[aliases[0]] = valor
        # QUEBRAS na fronteira: o texto pode vir de documento (VT/FF do Word) tanto
        # pela string solta quanto por este caminho, quando o modelo copia um trecho
        # de fonte para dentro do campo. Normalizar so no _texto_para_slide deixaria
        # metade do problema de pe - e a metade que aparece no slide.
        for campo in (_ALIAS_TEXTO[0], "titulo", "subtitulo", "antetitulo"):
            if isinstance(s.get(campo), str):
                s[campo] = _quebras(s[campo])
        for campo in (_ALIAS_BULLETS[0], _ALIAS_ITENS[0]):
            if isinstance(s.get(campo), list):
                s[campo] = [_quebras(x) if isinstance(x, str) else x
                            for x in s[campo]]
        tipo = _fold(s.get("tipo"))
        if tipo in _TIPOS_EXIGEM_CORPO:
            tem_corpo = bool(
                s.get("texto") or _lista(s.get("bullets")) or _lista(s.get("itens"))
            )
            if not tem_corpo:
                return slides, _diag_slide_mudo(metodo, idx, s)
    return slides, None


def _diag_entrada_vazia(metodo, nome_param, valor_raw):
    # Em vez de gerar um arquivo vazio em silencio, explica o que chegou.
    # (Mantido no chat: o conteudo e o proprio dado enviado, nada interno vaza.)
    log.warning(
        "gerador_nidum: %s recebeu 0 itens validos em '%s' (tipo=%s)",
        metodo, nome_param, type(valor_raw).__name__,
    )
    return (
        "DIAGNOSTICO " + metodo + ": recebi 0 itens validos em '" + nome_param
        + "' apos o parse, entao o arquivo ficaria vazio. Isso costuma significar "
        "que o conteudo nao foi enviado como lista de objetos. tipo_recebido="
        + type(valor_raw).__name__ + " | inicio=" + repr(valor_raw)[:400]
    )


def _erro_limpo(metodo):
    # v2.2.0: em vez de despejar o traceback no chat (vazava caminhos internos),
    # loga o traceback completo no servidor e devolve uma mensagem limpa com um
    # codigo curto para correlacionar no log.
    codigo = uuid.uuid4().hex[:8].upper()
    log.exception("gerador_nidum: %s falhou (codigo %s)", metodo, codigo)
    return (
        "Nao consegui gerar o arquivo desta vez (erro interno em " + metodo
        + "). Codigo do erro: " + codigo
        + " - informe este codigo a Tecnologia para localizar o detalhe nos logs."
    )


async def _salvar_e_linkar(data_bytes, filename, content_type, user_id):
    """Grava os bytes pelos modulos internos do Open WebUI (sem HTTP) e
    registra o arquivo, devolvendo link de download nativo.

    v2.2.0:
    - O upload roda em thread (asyncio.to_thread): Storage.upload_file e
      sincrono e, quando o Storage for rede (ex.: Cloudflare R2), bloquearia
      o event loop durante a transferencia.
    - TOLERANCIA A TAGS: provedores S3 sem suporte a x-amz-tagging no PutObject
      (caso documentado do Cloudflare R2) derrubam o upload com NotImplemented.
      A sequencia de tentativas agora inclui variantes SEM tags, entao o
      arquivo salva de qualquer forma (a tag de user_id e conveniencia, nao
      requisito).

    Nota de compatibilidade mantida da v2.1.0: Storage.upload_file e sincrono,
    mas Files.insert_new_file pode ser assincrono - usamos inspect.isawaitable
    para funcionar nas duas versoes do Open WebUI."""
    from open_webui.storage.provider import Storage
    from open_webui.models.files import Files, FileForm

    file_id = str(uuid.uuid4())
    stored_name = file_id + "_" + filename
    tags = {"OpenWebUI-User-Id": user_id} if user_id else {}

    def _upload():
        tentativas = []
        if tags:
            tentativas.append(
                ("bytesio+tags",
                 lambda: Storage.upload_file(io.BytesIO(data_bytes), stored_name, tags))
            )
        tentativas.append(
            ("bytesio",
             lambda: Storage.upload_file(io.BytesIO(data_bytes), stored_name, {}))
        )
        if tags:
            tentativas.append(
                ("raw+tags",
                 lambda: Storage.upload_file(data_bytes, stored_name, tags))
            )
        tentativas.append(
            ("raw",
             lambda: Storage.upload_file(data_bytes, stored_name, {}))
        )
        tentativas.append(
            ("raw-sem-param-tags",
             lambda: Storage.upload_file(data_bytes, stored_name))
        )
        last_err = None
        for rotulo, tentar in tentativas:
            try:
                return tentar()
            except Exception as e:
                last_err = e
                log.warning(
                    "gerador_nidum: upload variante '%s' falhou: %s", rotulo, e
                )
        raise RuntimeError("Falha ao salvar no Storage: " + str(last_err))

    result = await asyncio.to_thread(_upload)
    if isinstance(result, tuple) and len(result) >= 2:
        file_path = result[1]
    else:
        file_path = result
    if file_path is None:
        raise RuntimeError("Storage.upload_file devolveu caminho vazio.")

    meta = {"name": filename, "content_type": content_type, "size": len(data_bytes)}
    form = FileForm(id=file_id, filename=filename, path=file_path, meta=meta, data={})

    # Registro no banco. PRECISA de await na versao async.
    inserted = Files.insert_new_file(user_id, form)
    if inspect.isawaitable(inserted):
        inserted = await inserted
    if inserted is None:
        raise RuntimeError("Falha ao registrar o arquivo no banco (insert_new_file).")

    # LINK ABSOLUTO QUANDO POSSIVEL (defeito medido no laco agentico, 03/09/2026).
    # No pipe, esta string E a resposta: o caminho relativo chegava intacto ao usuario.
    # No laco agentico ela e INSUMO - o modelo reescreve o link na resposta dele, e ao
    # perder a barra inicial o navegador resolve relativo a /c/<chat_id>, gerando
    # /c/api/v1/files/... -> 404. URL absoluta e imune a essa resolucao.
    # WEBUI_URL vazia (default) -> mantem o caminho relativo, comportamento de sempre.
    download_path = "/api/v1/files/" + file_id + "/content"
    absoluto = False
    try:
        from open_webui.config import WEBUI_URL

        base = str(getattr(WEBUI_URL, "value", WEBUI_URL) or "").strip().rstrip("/")
        if base.startswith("http"):
            download_path = base + download_path
            absoluto = True
    except Exception:
        log.exception("gerador_nidum: falha ao resolver WEBUI_URL; link segue relativo")
    if not absoluto:
        # REGISTRO VISIVEL (pedido do Davi, 03/09/2026): configuracao vazia sem
        # registro e como a valve apontando para colecao apagada - ninguem descobre
        # ate quebrar. Com o anexo nativo abaixo isto deixa de ter consequencia.
        log.warning(
            "gerador_nidum: WEBUI_URL vazia -> o link de %s sai RELATIVO. Se o anexo "
            "nativo nao chegar, o modelo transcreve o caminho e pode gerar 404.",
            filename,
        )

    # ANEXO NATIVO primeiro: se ele chega, o endereco nunca passa pelo modelo.
    entregue = await _emitir_arquivo(
        _CTX_EMITTER.get(), _CTX_MESSAGE_ID.get(), file_id, filename,
        content_type, download_path,
    )
    if entregue:
        # Sem URL no texto - nao ha o que o modelo corromper. O card de download ja
        # esta na mensagem.
        return (
            "O arquivo " + filename + " foi anexado a esta resposta. O usuario ja tem "
            "o botao de download na tela - NAO escreva link nenhum, apenas diga que o "
            "arquivo esta pronto."
        )
    # FALLBACK DECLARADO: o evento nao chegou (sem emitter, ou falhou). O link volta,
    # e o texto DIZ que e fallback - para aparecer na tela em vez de virar 404 mudo.
    #
    # A FRASE "Link para download:" E CONTRATO COM O PIPE, NAO ENFEITE: o chatnd faz
    # `if "Link para download" in saida` (chatnd.py:5773) para decidir se oferece os
    # outros formatos. Mudar este texto desliga aquela oferta EM SILENCIO. O pipe
    # nunca passa __event_emitter__, entao ele cai sempre aqui e ve a string de
    # sempre - byte a byte. (Contrato por prosa entre dois componentes e fragil;
    # esta na fila para virar marcador explicito.)
    log.warning("gerador_nidum: anexo nativo indisponivel para %s; usando link", filename)
    return (
        "Link para download: [Clique aqui para baixar " + filename + "]("
        + download_path + ")"
        + "\n\n(Anexo nativo indisponivel nesta resposta - entregue por link. "
        "Reproduza o endereco acima EXATAMENTE como esta.)"
    )


class Tools:
    class Valves(BaseModel):
        # v2.2.0: OPENWEBUI_BASE_URL e OPENWEBUI_API_KEY removidas - nunca eram
        # usadas no codigo (campo de secret morto).
        # v2.3.0: ecossistema usado na nomenclatura quando o pipe nao informa uma sigla
        # valida (ou informa uma fora da lista fechada). Nunca deixa o nome sem prefixo.
        ECOSSISTEMA_PADRAO: str = "TEC"
        # v2.3.0: barra "Editar"/"Salvar HTML" no HTML e no deck gerados. Desligar aqui
        # se algum cliente nao quiser o editor embutido.
        EDITOR_HTML: bool = True

    def __init__(self):
        self.valves = self.Valves()

    # ------------------------------------------------------------------
    # Metodos publicos (expostos ao modelo). Erros internos: traceback vai
    # para o LOG do servidor; o chat recebe mensagem limpa com codigo de erro.
    # ------------------------------------------------------------------
    async def gerar_pptx(
        self, titulo: str, slides: list, marca: bool = True, __user__: dict = None,
        ecossistema: str = "", versao: int = 1, imagens: list = None,
        __event_emitter__=None, __message_id__: str = None
    ) -> str:
        """Gera uma apresentacao PowerPoint (.pptx) e devolve um link de download.

        :param titulo: titulo geral da apresentacao.
        :param slides: lista de slides. Cada slide e um dicionario com 'tipo',
            'titulo' e o conteudo do slide.

            O CONTEUDO VAI SEMPRE EM UM DESTES TRES CAMPOS, com estes nomes exatos:
              'texto'    - uma string. E o corpo do slide.
              'bullets'  - uma lista de strings.
              'itens'    - uma lista de {titulo, texto}, so em 'numerada' e 'cartoes'.
            Nao existe campo 'corpo', 'conteudo' nem 'descricao'. O corpo do slide
            se chama 'texto'.

            Nenhum slide de conteudo pode vir sem um desses tres: slide vazio sai
            mudo no arquivo e a ferramenta RECUSA a geracao com erro. A apresentacao
            precisa de ao menos 3 slides preenchidos.

            Os OITO tipos, e o que cada um espera:
              capa          - abertura. subtitulo (opcional).
              secao         - divisoria de tema, fundo colorido. 'cor' (opcional).
              conteudo      - titulo + 'texto' e/ou 'bullets' em fundo creme.
              destaque      - uma frase ou conceito forte em fundo colorido cheio.
                              'texto' + 'cor'.
              divisao       - titulo num bloco de cor a esquerda, 'texto'/'bullets'
                              a direita. Defina 'cor'.
              numerada      - etapas ou passos com numeros grandes. Preencha
                              'itens' com uma lista de {titulo, texto}.
              cartoes       - 2 a 4 cartoes coloridos lado a lado (valores,
                              pilares). Preencha 'itens' com {titulo, texto}.
              encerramento  - fecho.

            Campos aceitos: tipo, titulo, subtitulo, texto (string), bullets
            (lista de strings), itens (lista de {titulo, texto} - so em
            'numerada' e 'cartoes'), cor ('verde'|'azul'|'terracota'|'preto'),
            imagem (marcador tipo 'IMAGEM_1').

            VARIE os layouts - nao use so 'conteudo', fica monotono. Numa
            apresentacao tipica alterne os tipos (ex.: capa, conteudo, destaque,
            cartoes, divisao, numerada, secao, encerramento) e varie a 'cor'
            entre slides coloridos vizinhos. Para pilares, valores ou categorias
            prefira 'cartoes'; para etapas ou passos, 'numerada'; para uma
            afirmacao de impacto, 'destaque'. Prefira bullets curtos a
            paragrafos longos.
        :param marca: aplica a identidade visual da Nidum (padrao True).
        :param ecossistema: sigla do ecossistema para a nomenclatura oficial (opcional).
        :param versao: numero de versao para o nome do arquivo (padrao 1).
        :param imagens: imagens ANEXADAS PELO USUARIO (data-URL base64 ou bytes), na
            ordem dos marcadores IMAGEM_1, IMAGEM_2... Vem do pipe por parametro,
            nunca de um modelo.
        :return: um link markdown pronto, ja no formato "[Clique aqui para baixar
            <nome>](<url>)". REPRODUZA ESSE LINK EXATAMENTE COMO RECEBIDO - copie e
            cole, sem reescrever, sem encurtar e sem remover a barra inicial da URL.
            Um caractere a menos quebra o download (o navegador resolve o caminho
            relativo a pagina do chat e devolve 404).
        """
        _ctx_entrega(__event_emitter__, __message_id__)
        try:
            return await self._pptx(titulo, slides, marca, __user__, ecossistema,
                                    versao, imagens)
        except Exception:
            return _erro_limpo("gerar_pptx")

    async def gerar_xlsx(
        self, titulo: str, planilhas: list, marca: bool = True, __user__: dict = None,
        ecossistema: str = "", versao: int = 1,
        __event_emitter__=None, __message_id__: str = None
    ) -> str:
        """Gera uma planilha Excel (.xlsx) e devolve um link de download.

        :param titulo: titulo geral.
        :param planilhas: lista de abas. Cada aba e um dicionario com:
            nome, cabecalhos (lista), linhas (lista de listas).
        :param marca: aplica a identidade visual da Nidum (padrao True).
        :param ecossistema: sigla do ecossistema para a nomenclatura oficial (opcional).
        :param versao: numero de versao para o nome do arquivo (padrao 1).
        :return: um link markdown pronto, ja no formato "[Clique aqui para baixar
            <nome>](<url>)". REPRODUZA ESSE LINK EXATAMENTE COMO RECEBIDO - copie e
            cole, sem reescrever, sem encurtar e sem remover a barra inicial da URL.
            Um caractere a menos quebra o download (o navegador resolve o caminho
            relativo a pagina do chat e devolve 404).
        """
        _ctx_entrega(__event_emitter__, __message_id__)
        try:
            return await self._xlsx(titulo, planilhas, marca, __user__, ecossistema, versao)
        except Exception:
            return _erro_limpo("gerar_xlsx")

    async def gerar_docx(
        self, titulo: str, secoes: list, marca: bool = True, __user__: dict = None,
        ecossistema: str = "", versao: int = 1, imagens: list = None,
        __event_emitter__=None, __message_id__: str = None
    ) -> str:
        """Gera um documento Word (.docx) e devolve um link de download.

        :param titulo: titulo do documento.
        :param secoes: lista de secoes. Cada secao e um dicionario com:
            heading, paragrafos (opcional, lista), bullets (opcional, lista),
            imagem (opcional, marcador tipo 'IMAGEM_1').
        :param marca: aplica a identidade visual da Nidum (padrao True).
        :param ecossistema: sigla do ecossistema para a nomenclatura oficial (opcional).
        :param versao: numero de versao para o nome do arquivo (padrao 1).
        :param imagens: imagens ANEXADAS PELO USUARIO (data-URL base64 ou bytes), na
            ordem dos marcadores IMAGEM_1, IMAGEM_2...
        :return: um link markdown pronto, ja no formato "[Clique aqui para baixar
            <nome>](<url>)". REPRODUZA ESSE LINK EXATAMENTE COMO RECEBIDO - copie e
            cole, sem reescrever, sem encurtar e sem remover a barra inicial da URL.
            Um caractere a menos quebra o download (o navegador resolve o caminho
            relativo a pagina do chat e devolve 404).
        """
        _ctx_entrega(__event_emitter__, __message_id__)
        try:
            return await self._docx(titulo, secoes, marca, __user__, ecossistema,
                                    versao, imagens)
        except Exception:
            return _erro_limpo("gerar_docx")

    async def gerar_pdf(
        self, titulo: str, secoes: list, marca: bool = True, __user__: dict = None,
        ecossistema: str = "", versao: int = 1, imagens: list = None,
        __event_emitter__=None, __message_id__: str = None
    ) -> str:
        """Gera um documento PDF e devolve um link de download.

        :param titulo: titulo do documento.
        :param secoes: lista de secoes. Cada secao e um dicionario com:
            heading, paragrafos (opcional, lista), bullets (opcional, lista),
            tabela (opcional, lista de listas), imagem (opcional, 'IMAGEM_1').
        :param marca: aplica a identidade visual da Nidum (padrao True).
        :param ecossistema: sigla do ecossistema para a nomenclatura oficial (opcional).
        :param versao: numero de versao para o nome do arquivo (padrao 1).
        :param imagens: imagens ANEXADAS PELO USUARIO (data-URL base64 ou bytes), na
            ordem dos marcadores IMAGEM_1, IMAGEM_2...
        :return: um link markdown pronto, ja no formato "[Clique aqui para baixar
            <nome>](<url>)". REPRODUZA ESSE LINK EXATAMENTE COMO RECEBIDO - copie e
            cole, sem reescrever, sem encurtar e sem remover a barra inicial da URL.
            Um caractere a menos quebra o download (o navegador resolve o caminho
            relativo a pagina do chat e devolve 404).
        """
        _ctx_entrega(__event_emitter__, __message_id__)
        try:
            return await self._pdf(titulo, secoes, marca, __user__, ecossistema,
                                   versao, imagens)
        except Exception:
            return _erro_limpo("gerar_pdf")

    # Content-type por extensao para o MODO PRESERVACAO (gerar_codigo). So familia
    # texto/codigo - o que "e o proprio codigo-fonte". Binario (xlsx/pptx) NAO entra:
    # precisa de parser, e outra fatia (documentada, nao construida).
    _CT_CODIGO = {
        "html": "text/html", "htm": "text/html", "css": "text/css",
        "js": "text/javascript", "json": "application/json", "xml": "application/xml",
        "md": "text/markdown", "txt": "text/plain", "csv": "text/csv",
    }

    async def gerar_codigo(
        self, titulo: str, conteudo: str, ext: str = "html", __user__: dict = None,
        ecossistema: str = "", versao: int = 1,
        __event_emitter__=None, __message_id__: str = None
    ) -> str:
        """MODO PRESERVACAO: grava codigo-fonte VERBATIM, sem marca e sem editor.

        Existe para EDITAR arquivo do usuario (ex.: um app HTML com <script> e handlers)
        preservando o comportamento. Ao contrario do gerar_html, NAO injeta a identidade
        Nidum (repintaria o CSS do app) nem a barra de edicao (o contenteditable no body
        briga com os campos de formulario, e o app ja tem os proprios controles). O que
        entra e o que sai - byte a byte, sem reescrever nada.

        :param conteudo: o arquivo INTEIRO ja editado (string), tal como deve ser salvo.
        :param ext: extensao da familia texto/codigo (html/css/js/json/xml/md/txt/csv).
        :return: um link markdown pronto, ja no formato "[Clique aqui para baixar
            <nome>](<url>)". REPRODUZA ESSE LINK EXATAMENTE COMO RECEBIDO - copie e
            cole, sem reescrever, sem encurtar e sem remover a barra inicial da URL.
            Um caractere a menos quebra o download (o navegador resolve o caminho
            relativo a pagina do chat e devolve 404).
        """
        _ctx_entrega(__event_emitter__, __message_id__)
        try:
            # NAO passa por _coerce: ele faz json.loads e um arquivo .json literal viraria
            # um dict reserializado (aspas trocadas) - o oposto de "verbatim". Codigo e
            # string crua, ponto.
            c = conteudo if isinstance(conteudo, str) else str(conteudo or "")
            if not c.strip():
                return _diag_entrada_vazia("gerar_codigo", "conteudo", conteudo)
            e = str(ext or "html").lstrip(".").lower()
            ct = self._CT_CODIGO.get(e)
            if ct is None:
                # Fora da familia texto -> nao e para ca (binario precisa de parser).
                log.warning("gerador_nidum: gerar_codigo com ext %r fora da familia; "
                            "tratando como txt", e)
                e, ct = "txt", "text/plain"
            nome = _nome_padrao(titulo, ecossistema, e, versao,
                                self.valves.ECOSSISTEMA_PADRAO)
            data = c.encode("utf-8")
            link = await _salvar_e_linkar(data, nome, ct, _get_user_id(__user__))
            return "Arquivo gerado com sucesso. " + link
        except Exception:
            return _erro_limpo("gerar_codigo")

    async def gerar_html(
        self, titulo: str, html: str, __user__: dict = None,
        ecossistema: str = "", versao: int = 1, imagens: list = None,
        __event_emitter__=None, __message_id__: str = None
    ) -> str:
        """Gera um arquivo HTML (.html) e devolve um link de download.

        :param titulo: titulo/nome do arquivo.
        :param html: documento HTML completo (string), pronto para abrir no navegador.
        :param ecossistema: sigla do ecossistema para a nomenclatura oficial (opcional).
        :param versao: numero de versao para o nome do arquivo (padrao 1).
        :param imagens: imagens ANEXADAS PELO USUARIO. O modelo posiciona o marcador
            (IMAGEM_1) no HTML e aqui ele vira a imagem embutida em base64.
        :return: um link markdown pronto, ja no formato "[Clique aqui para baixar
            <nome>](<url>)". REPRODUZA ESSE LINK EXATAMENTE COMO RECEBIDO - copie e
            cole, sem reescrever, sem encurtar e sem remover a barra inicial da URL.
            Um caractere a menos quebra o download (o navegador resolve o caminho
            relativo a pagina do chat e devolve 404).
        """
        _ctx_entrega(__event_emitter__, __message_id__)
        try:
            conteudo = _coerce(html)
            conteudo = conteudo if isinstance(conteudo, str) else str(conteudo or "")
            conteudo = conteudo.strip()
            if not conteudo:
                return _diag_entrada_vazia("gerar_html", "html", html)

            # Nome calculado ANTES da montagem: o bloco do editor precisa dele (e o nome
            # sugerido no download do "Salvar HTML").
            nome = _nome_padrao(titulo, ecossistema, "html", versao,
                                self.valves.ECOSSISTEMA_PADRAO)
            com_editor = self.valves.EDITOR_HTML
            imgs = _normalizar_imagens(imagens)

            def _montar():
                c = conteudo
                # Marcadores -> imagens reais ANTES de embrulhar/injetar marca, para a
                # imagem entrar no corpo do documento como qualquer outro conteudo.
                if imgs:
                    c = _inserir_imagens_html(c, imgs)
                # Se vier so um fragmento, embrulha num documento HTML minimo. O titulo e
                # ESCAPADO: '<' ou '&' no titulo quebrariam a tag <title> (bug 2.5).
                if "<html" not in c.lower():
                    c = (
                        '<!DOCTYPE html>\n<html lang="pt-br"><head><meta charset="utf-8">'
                        '<meta name="viewport" content="width=device-width, initial-scale=1">'
                        "<title>" + _esc(titulo or "Documento") + "</title></head><body>\n"
                        + c + "\n</body></html>"
                    )
                # _injetar_marca_html le fontes/logo do disco (IO) - por isso
                # roda aqui dentro da thread.
                c = _injetar_marca_html(c)
                if com_editor:
                    c = _injetar_editor(c, nome)
                return c.encode("utf-8")

            data = await asyncio.to_thread(_montar)
            link = await _salvar_e_linkar(
                data, nome, "text/html", _get_user_id(__user__)
            )
            return "Arquivo gerado com sucesso. " + link
        except Exception:
            return _erro_limpo("gerar_html")

    async def gerar_apresentacao_html(
        self, titulo: str, slides: list, __user__: dict = None,
        ecossistema: str = "", versao: int = 1, imagens: list = None,
        __event_emitter__=None, __message_id__: str = None
    ) -> str:
        """Gera uma APRESENTACAO em HTML navegavel (deck) com a identidade Nidum.

        Deck autocontido: 1 slide por vez, navegacao por setas/teclado/dots,
        cantos arredondados, transicoes, contraste correto e fonte embutida.
        :param slides: mesma estrutura do gerar_pptx (lista de slides com tipo,
            titulo, subtitulo, texto, bullets, cor, itens, imagem).
        :param ecossistema: sigla do ecossistema para a nomenclatura oficial (opcional).
        :param versao: numero de versao para o nome do arquivo (padrao 1).
        :param imagens: imagens ANEXADAS PELO USUARIO, na ordem dos marcadores.
        :return: um link markdown pronto, ja no formato "[Clique aqui para baixar
            <nome>](<url>)". REPRODUZA ESSE LINK EXATAMENTE COMO RECEBIDO - copie e
            cole, sem reescrever, sem encurtar e sem remover a barra inicial da URL.
            Um caractere a menos quebra o download (o navegador resolve o caminho
            relativo a pagina do chat e devolve 404).
        """
        _ctx_entrega(__event_emitter__, __message_id__)
        try:
            return await self._apresentacao_html(
                titulo, slides, __user__, ecossistema, versao, imagens
            )
        except Exception:
            return _erro_limpo("gerar_apresentacao_html")

    # ------------------------------------------------------------------
    # Implementacoes (privadas - nao expostas ao modelo).
    # v2.2.0: a montagem pesada roda em asyncio.to_thread para NAO travar o
    # event loop do Open WebUI (um render grande congelava todos os usuarios).
    # ------------------------------------------------------------------
    async def _apresentacao_html(self, titulo, slides, __user__,
                                 ecossistema="", versao=1, imagens=None):
        raw = slides
        slides = _itens_loose(slides, _texto_para_slide)
        if not slides:
            return _diag_entrada_vazia("gerar_apresentacao_html", "slides", raw)
        slides, erro = _normalizar_corpo_slides(slides, "gerar_apresentacao_html")
        if erro:
            return erro
        imgs = _normalizar_imagens(imagens)

        # Nome ANTES da montagem (o editor precisa dele para o download).
        nome = _nome_padrao(titulo, ecossistema, "html", versao,
                            self.valves.ECOSSISTEMA_PADRAO)
        com_editor = self.valves.EDITOR_HTML

        def _montar():
            faces = "".join(
                [
                    _font_face("MaximaNouva-Thin.ttf", 300),
                    _font_face("MaximaNouva-Regular.ttf", 400),
                    _font_face("MaximaNouva-SemiBold.ttf", 600),
                    _font_face("MaximaNouva-Bold.ttf", 700),
                    _font_face("MaximaNouva-ExtraBold.ttf", 800),
                    _font_face("MaximaNouva-Italic.ttf", 400, "italic"),
                    _font_face("Ibrand.ttf", "100 900", "normal", "Ibrand"),
                ]
            )
            logo_t = _logo_b64("terracota")
            logo_a = _logo_b64("areia")
            mapa = {
                "verde": "#515E52", "azul": "#4F7187", "terracota": "#9A4A2E",
                "preto": "#1F1E1B", "creme": "#E5E0D5",
            }
            cores_secao = ["#515E52", "#4F7187", "#9A4A2E"]
            cores_cartao = ["#4F7187", "#515E52", "#9A4A2E", "#1F1E1B"]
            sec = [0]

            partes = []
            for s in slides:
                tipo = _fold(s.get("tipo")) or "conteudo"
                partes.append(
                    _slide_html(
                        s, tipo, mapa, cores_secao, cores_cartao, sec, logo_t, logo_a
                    )
                )
                # Imagem do usuario: slide PROPRIO logo apos, igual ao pptx (os layouts
                # do deck sao composicoes fechadas; a foto em slide proprio aparece
                # grande e nao briga com o texto). object-fit:contain = nunca distorce.
                img_s = _imagem_do_item(imgs, s)
                if img_s is not None:
                    uri = _img_data_uri(img_s)
                    if uri:
                        leg = _esc(s.get("titulo"))
                        partes.append(
                            "<section class='slide'>"
                            + (("<h2>" + leg + "</h2>") if leg else "")
                            + "<img src='" + uri + "' alt='' style='display:block;"
                            "margin:18px auto 0;max-width:88%;max-height:66vh;"
                            "width:auto;height:auto;object-fit:contain;"
                            "border-radius:12px'>"
                            + (("<img class='logo' src='" + logo_t + "'>")
                               if logo_t else "")
                            + "</section>"
                        )
            deck = "".join(partes)
            # Os dots contam os slides RENDERIZADOS (as imagens acrescentam slides).
            dots = "".join(
                "<span class='dot' onclick='go(" + str(k) + ")'></span>"
                for k in range(len(partes))
            )
            nav = (
                "<div class='nav'><button onclick='go(i-1)'>&#8249;</button>"
                "<div class='dots'>" + dots + "</div>"
                "<button onclick='go(i+1)'>&#8250;</button>"
                "<span class='count' id='cnt'></span></div>"
            )
            html = (
                "<!DOCTYPE html><html lang=\"pt-br\"><head><meta charset=\"utf-8\">"
                "<meta name=\"viewport\" content=\"width=device-width,initial-scale=1\">"
                "<title>" + _esc(titulo or "Apresentacao Nidum") + "</title>"
                "<style>" + faces + DECK_CSS + "</style></head><body>"
                "<div class='deck'>" + deck + "</div>" + nav + DECK_JS + "</body></html>"
            )
            if com_editor:
                html = _injetar_editor(html, nome)
            return html.encode("utf-8")

        data = await asyncio.to_thread(_montar)
        link = await _salvar_e_linkar(
            data, nome, "text/html", _get_user_id(__user__)
        )
        return "Arquivo gerado com sucesso. " + link

    async def _pptx(self, titulo, slides, marca, __user__,
                    ecossistema="", versao=1, imagens=None):
        raw_slides = slides
        slides = _itens_loose(slides, _texto_para_slide)
        if not slides:
            return _diag_entrada_vazia("gerar_pptx", "slides", raw_slides)
        slides, erro = _normalizar_corpo_slides(slides, "gerar_pptx")
        if erro:
            return erro
        imgs = _normalizar_imagens(imagens)

        def _montar():
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.enum.shapes import MSO_SHAPE

            def col(h):
                return RGBColor(*_hex_to_rgb(h))

            terracota = col(NIDUM_TERRACOTA)
            verde = col(NIDUM_VERDE)
            azul = col(NIDUM_AZUL)
            cinza = col(NIDUM_CINZA)
            preto = col(NIDUM_PRETO)
            creme = col(NIDUM_CREME)
            branco = col(NIDUM_BRANCO)
            cores_secao = [verde, azul, terracota]
            cores_cartao = [azul, verde, terracota, preto]
            mapa_cor = {"verde": verde, "azul": azul, "terracota": terracota,
                        "preto": preto, "creme": creme}
            sec_idx = 0

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            blank = prs.slide_layouts[6]
            SW = prs.slide_width
            SH = prs.slide_height

            def add_fundo(slide, cor):
                shp = slide.shapes.add_shape(1, 0, 0, SW, SH)
                shp.fill.solid()
                shp.fill.fore_color.rgb = cor
                shp.line.fill.background()
                shp.shadow.inherit = False
                slide.shapes._spTree.remove(shp._element)
                slide.shapes._spTree.insert(2, shp._element)
                return shp

            def add_caixa(slide, left, top, width, height, anchor=None):
                tb = slide.shapes.add_textbox(left, top, width, height)
                tf = tb.text_frame
                tf.word_wrap = True
                if anchor is not None:
                    tf.vertical_anchor = anchor
                return tf

            def estilo(p, size, cor, bold=False, upper=False, italic=False, align=None):
                if upper and p.text:
                    p.text = p.text.upper()
                # Titulos (bold) usam a Ibrand (com peso proprio); corpo usa Maxima.
                fonte = NIDUM_FONT_LOGO if bold else NIDUM_FONT
                negr = bold and fonte == NIDUM_FONT
                p.font.name = fonte
                p.font.size = Pt(size)
                p.font.bold = negr
                p.font.italic = italic
                p.font.color.rgb = cor
                if align is not None:
                    p.alignment = align
                for r in p.runs:
                    r.font.name = fonte
                    r.font.size = Pt(size)
                    r.font.bold = negr
                    r.font.italic = italic
                    r.font.color.rgb = cor

            def add_logo(slide, cor_logo, left, top, width):
                p = _logo_path(cor_logo)
                if not p:
                    return
                try:
                    slide.shapes.add_picture(p, left, top, width=width)
                except Exception:
                    log.exception("gerador_nidum: falha ao inserir logo no pptx")

            def add_bloco(slide, left, top, w, h, cor):
                sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
                sh.fill.solid()
                sh.fill.fore_color.rgb = cor
                sh.line.fill.background()
                sh.shadow.inherit = False
                return sh

            def add_card(slide, left, top, w, h, cor, titulo_c, texto_c):
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h
                )
                card.fill.solid()
                card.fill.fore_color.rgb = cor
                card.line.fill.background()
                card.shadow.inherit = False
                tf = card.text_frame
                tf.word_wrap = True
                tf.margin_left = Inches(0.22)
                tf.margin_right = Inches(0.22)
                tf.margin_top = Inches(0.18)
                p = tf.paragraphs[0]
                p.text = titulo_c or ""
                estilo(p, 15, creme, bold=True)
                if texto_c:
                    p2 = tf.add_paragraph()
                    p2.text = texto_c
                    estilo(p2, 11, creme)
                return card

            def cor_de(nome, padrao):
                return mapa_cor.get(str(nome).lower(), padrao) if nome else padrao

            def add_slide_imagem(img, legenda):
                # A imagem do usuario ganha o PROPRIO slide, logo apos o slide que a
                # posicionou. Motivo: os layouts (capa, cartoes, numerada, destaque...)
                # sao composicoes fechadas - encaixar uma foto dentro deles colidiria com
                # texto e quebraria a regua de marca. Em slide proprio a imagem aparece
                # grande, centralizada, com margem, e o deck continua limpo.
                sl = prs.slides.add_slide(blank)
                if marca:
                    add_fundo(sl, creme)
                if legenda:
                    tfl = add_caixa(sl, Inches(0.9), Inches(0.45), Inches(11.5),
                                    Inches(0.8))
                    pl = tfl.paragraphs[0]
                    pl.text = str(legenda)
                    estilo(pl, 18, verde if marca else preto, bold=True,
                           align=PP_ALIGN.CENTER)
                topo = Inches(1.35) if legenda else Inches(0.8)
                caixa_w = SW - Inches(2.0)
                caixa_h = Inches(6.3) - topo
                try:
                    # add_picture sem width/height entra no TAMANHO NATIVO; dai medimos e
                    # reescalamos com fator unico (nunca distorce).
                    pic = sl.shapes.add_picture(io.BytesIO(img["bytes"]), 0, topo)
                    w, h = _encaixar(pic.width, pic.height, caixa_w, caixa_h)
                    pic.width = int(w)
                    pic.height = int(h)
                    pic.left = int((SW - pic.width) / 2)
                    pic.top = int(topo + (caixa_h - pic.height) / 2)
                except Exception:
                    log.exception(
                        "gerador_nidum: falha ao inserir imagem do usuario no pptx"
                    )
                    return
                if marca:
                    add_logo(sl, "terracota", Inches(11.1), Inches(6.8), Inches(1.6))

            for s in slides:
                tipo = _fold(s.get("tipo")) or "conteudo"
                slide = prs.slides.add_slide(blank)
                cor_titulo = verde if marca else preto
                cor_corpo = preto

                if tipo == "capa":
                    if marca:
                        add_fundo(slide, creme)
                        add_logo(slide, "terracota", Inches(3.67), Inches(1.3), Inches(6.0))
                    tf = add_caixa(
                        slide, Inches(1.0), Inches(4.4), Inches(11.33), Inches(2.0),
                        MSO_ANCHOR.TOP,
                    )
                    p = tf.paragraphs[0]
                    p.text = s.get("titulo") or titulo or ""
                    estilo(p, 40, cor_titulo, bold=True, align=PP_ALIGN.CENTER)
                    if s.get("subtitulo"):
                        p2 = tf.add_paragraph()
                        p2.text = s["subtitulo"]
                        estilo(p2, 18, terracota if marca else preto, align=PP_ALIGN.CENTER)

                elif tipo == "secao":
                    cor_sec = cores_secao[sec_idx % len(cores_secao)]
                    sec_idx += 1
                    cor_sec = cor_de(s.get("cor"), cor_sec)
                    if marca:
                        add_fundo(slide, cor_sec)
                    tf = add_caixa(
                        slide, Inches(1.0), Inches(2.6), Inches(11.33), Inches(2.3),
                        MSO_ANCHOR.MIDDLE,
                    )
                    if s.get("subtitulo"):
                        pa = tf.paragraphs[0]
                        pa.text = s["subtitulo"]
                        estilo(pa, 16, creme if marca else preto, upper=True)
                        p = tf.add_paragraph()
                    else:
                        p = tf.paragraphs[0]
                    p.text = s.get("titulo", "")
                    estilo(p, 40, creme if marca else preto, bold=True)
                    if marca:
                        add_logo(slide, "areia", Inches(0.7), Inches(6.6), Inches(1.8))

                elif tipo == "destaque":
                    cor_d = cor_de(s.get("cor"), terracota)
                    if marca:
                        add_fundo(slide, cor_d)
                    tf = add_caixa(
                        slide, Inches(1.2), Inches(2.2), Inches(10.9), Inches(3.1),
                        MSO_ANCHOR.MIDDLE,
                    )
                    if s.get("subtitulo"):
                        pa = tf.paragraphs[0]
                        pa.text = s["subtitulo"]
                        estilo(pa, 16, creme if marca else preto, upper=True)
                        p = tf.add_paragraph()
                    else:
                        p = tf.paragraphs[0]
                    p.text = s.get("titulo", "")
                    estilo(p, 34, creme if marca else preto, bold=True)
                    if s.get("texto"):
                        p2 = tf.add_paragraph()
                        p2.text = s["texto"]
                        estilo(p2, 16, creme if marca else preto)
                    if marca:
                        add_logo(slide, "areia", Inches(11.0), Inches(6.8), Inches(1.6))

                elif tipo == "divisao":
                    cor_b = cor_de(s.get("cor"), verde)
                    if marca:
                        add_fundo(slide, creme)
                        add_bloco(slide, 0, 0, Inches(5.4), SH, cor_b)
                    tf_l = add_caixa(
                        slide, Inches(0.7), Inches(0.9), Inches(4.2), Inches(5.6),
                        MSO_ANCHOR.MIDDLE,
                    )
                    pl = tf_l.paragraphs[0]
                    pl.text = s.get("titulo", "")
                    estilo(pl, 28, creme if marca else preto, bold=True)
                    tf_r = add_caixa(slide, Inches(5.9), Inches(0.9), Inches(6.7), Inches(5.6))
                    first = True
                    if s.get("subtitulo"):
                        pa = tf_r.paragraphs[0]
                        pa.text = s["subtitulo"]
                        estilo(pa, 13, preto, upper=True)
                        first = False
                    if s.get("texto"):
                        p = tf_r.paragraphs[0] if first else tf_r.add_paragraph()
                        first = False
                        p.text = s["texto"]
                        estilo(p, 15, preto)
                    for b in _lista(s.get("bullets")):
                        p = tf_r.paragraphs[0] if first else tf_r.add_paragraph()
                        first = False
                        p.text = str(b)
                        estilo(p, 15, preto)

                elif tipo == "numerada":
                    if marca:
                        add_fundo(slide, creme)
                    tf_t = add_caixa(slide, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.3))
                    if s.get("subtitulo"):
                        pa = tf_t.paragraphs[0]
                        pa.text = s["subtitulo"]
                        estilo(pa, 14, preto, upper=True)
                        pt = tf_t.add_paragraph()
                    else:
                        pt = tf_t.paragraphs[0]
                    pt.text = s.get("titulo", "")
                    estilo(pt, 28, verde if marca else preto, bold=True)
                    itens = _lista_de_dicts(s.get("itens"))
                    norm = []
                    for it in itens:
                        norm.append(
                            (it.get("titulo") or "",
                             it.get("texto") or " ".join(_lista(it.get("bullets"))))
                        )
                    if not norm:
                        norm = [("", str(b)) for b in _lista(s.get("bullets"))]
                    for idx, (t_i, txt_i) in enumerate(norm[:6]):
                        y = Inches(2.0 + idx * 0.82)
                        tfn = add_caixa(slide, Inches(0.9), y, Inches(0.9), Inches(0.8))
                        pn = tfn.paragraphs[0]
                        pn.text = str(idx + 1)
                        estilo(pn, 30, terracota, bold=True)
                        tfx = add_caixa(slide, Inches(1.9), y, Inches(10.5), Inches(0.8),
                                        MSO_ANCHOR.MIDDLE)
                        px = tfx.paragraphs[0]
                        px.text = (t_i + " - " + txt_i).strip(" -") if t_i else txt_i
                        estilo(px, 15, preto)

                elif tipo == "cartoes":
                    if marca:
                        add_fundo(slide, creme)
                    tf_t = add_caixa(slide, Inches(0.9), Inches(0.55), Inches(11.5), Inches(1.1))
                    if s.get("subtitulo"):
                        pa = tf_t.paragraphs[0]
                        pa.text = s["subtitulo"]
                        estilo(pa, 14, preto, upper=True)
                        pt = tf_t.add_paragraph()
                    else:
                        pt = tf_t.paragraphs[0]
                    pt.text = s.get("titulo", "")
                    estilo(pt, 26, verde if marca else preto, bold=True)
                    itens = _lista_de_dicts(s.get("itens"))
                    if not itens:
                        itens = [{"titulo": "", "texto": str(b)} for b in _lista(s.get("bullets"))]
                    itens = itens[:4]
                    xs = [Inches(0.9), Inches(6.85)]
                    ys = [Inches(1.9), Inches(4.35)]
                    for idx, it in enumerate(itens):
                        add_card(
                            slide, xs[idx % 2], ys[idx // 2], Inches(5.6), Inches(2.2),
                            cores_cartao[idx % len(cores_cartao)],
                            it.get("titulo") or "", it.get("texto") or "",
                        )

                elif tipo == "encerramento":
                    if marca:
                        add_fundo(slide, creme)
                        add_logo(slide, "terracota", Inches(4.67), Inches(2.4), Inches(4.0))
                    tf = add_caixa(
                        slide, Inches(1.0), Inches(4.6), Inches(11.33), Inches(1.2),
                        MSO_ANCHOR.TOP,
                    )
                    p = tf.paragraphs[0]
                    p.text = s.get("titulo") or "Fazer da casa um ninho."
                    estilo(p, 22, cor_titulo, align=PP_ALIGN.CENTER)
                    if s.get("texto"):
                        p2 = tf.add_paragraph()
                        p2.text = s["texto"]
                        estilo(p2, 14, preto, align=PP_ALIGN.CENTER)

                else:
                    if marca:
                        add_fundo(slide, creme)
                    tf_t = add_caixa(slide, Inches(0.9), Inches(0.7), Inches(11.5), Inches(1.4))
                    if s.get("subtitulo"):
                        pa = tf_t.paragraphs[0]
                        pa.text = s["subtitulo"]
                        estilo(pa, 14, preto, upper=True)
                        pt = tf_t.add_paragraph()
                    else:
                        pt = tf_t.paragraphs[0]
                    pt.text = s.get("titulo", "")
                    estilo(pt, 30, cor_titulo, bold=True)

                    tf_b = add_caixa(slide, Inches(0.9), Inches(2.1), Inches(11.5), Inches(4.5))
                    first = True
                    if s.get("texto"):
                        p = tf_b.paragraphs[0]
                        p.text = s["texto"]
                        estilo(p, 16, cor_corpo)
                        first = False
                    for b in _lista(s.get("bullets")):
                        p = tf_b.paragraphs[0] if first else tf_b.add_paragraph()
                        first = False
                        p.text = str(b)
                        estilo(p, 16, cor_corpo)
                    if marca:
                        add_logo(slide, "terracota", Inches(11.1), Inches(6.8), Inches(1.6))

                # Imagem que o usuario anexou e o modelo posicionou NESTE slide.
                img_s = _imagem_do_item(imgs, s)
                if img_s is not None:
                    add_slide_imagem(img_s, s.get("titulo") or "")

            buf = io.BytesIO()
            prs.save(buf)
            return buf.getvalue()

        data = await asyncio.to_thread(_montar)
        nome = _nome_padrao(titulo, ecossistema, "pptx", versao, self.valves.ECOSSISTEMA_PADRAO)
        ct = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        link = await _salvar_e_linkar(data, nome, ct, _get_user_id(__user__))
        return "Arquivo gerado com sucesso. " + link

    async def _xlsx(self, titulo, planilhas, marca, __user__,
                    ecossistema="", versao=1):
        raw_planilhas = planilhas
        planilhas = _lista_de_dicts(planilhas)
        if not planilhas:
            return _diag_entrada_vazia("gerar_xlsx", "planilhas", raw_planilhas)

        def _montar():
            from openpyxl import Workbook
            from openpyxl.styles import PatternFill, Font, Alignment, Border, Side
            from openpyxl.utils import get_column_letter

            wb = Workbook()
            wb.remove(wb.active)

            f_titulo = Font(name=NIDUM_FONT_LOGO, size=18, bold=False, color=NIDUM_VERDE)
            f_header = Font(name=NIDUM_FONT, size=11, bold=True, color=NIDUM_CREME)
            f_body = Font(name=NIDUM_FONT, size=11, color=NIDUM_PRETO)
            fill_header = PatternFill("solid", fgColor=NIDUM_VERDE)
            fill_alt = PatternFill("solid", fgColor=NIDUM_CREME_ALT)
            linha_side = Side(style="thin", color=NIDUM_CINZA)
            borda_b = Border(bottom=linha_side)
            al_h = Alignment(horizontal="left", vertical="center")
            al_b = Alignment(horizontal="left", vertical="center", wrap_text=True)

            for pl in planilhas:
                ws = wb.create_sheet(title=(pl.get("nome") or "Planilha")[:31])
                if marca:
                    ws.sheet_properties.tabColor = NIDUM_VERDE
                cabec = _lista(pl.get("cabecalhos"))
                linhas = [_lista(r) for r in _lista(pl.get("linhas"))]
                ncols = max([len(cabec)] + [len(r) for r in linhas] + [1])

                r = 1
                if marca:
                    ws.merge_cells(
                        start_row=1, start_column=1, end_row=1, end_column=ncols
                    )
                    ct = ws.cell(row=1, column=1, value=(pl.get("nome") or titulo or "Planilha"))
                    ct.font = f_titulo
                    ct.alignment = al_h
                    ws.row_dimensions[1].height = 28
                    r = 2

                header_row = None
                if cabec:
                    header_row = r
                    for j, val in enumerate(cabec, start=1):
                        cell = ws.cell(row=r, column=j, value=val)
                        cell.alignment = al_h
                        if marca:
                            cell.fill = fill_header
                            cell.font = f_header
                    ws.row_dimensions[r].height = 22
                    r += 1

                data_ini = r
                for ld in linhas:
                    for j in range(1, ncols + 1):
                        val = ld[j - 1] if j - 1 < len(ld) else None
                        cell = ws.cell(row=r, column=j, value=val)
                        cell.font = f_body
                        cell.alignment = al_b
                        cell.border = borda_b
                        if marca and ((r - data_ini) % 2 == 1):
                            cell.fill = fill_alt
                    r += 1

                for j in range(1, ncols + 1):
                    maxlen = len(str(cabec[j - 1])) if cabec and j - 1 < len(cabec) else 0
                    for ld in linhas:
                        if j - 1 < len(ld) and ld[j - 1] is not None:
                            maxlen = max(maxlen, len(str(ld[j - 1])))
                    ws.column_dimensions[get_column_letter(j)].width = min(
                        max(maxlen + 3, 12), 44
                    )

                if header_row:
                    ws.freeze_panes = ws.cell(row=header_row + 1, column=1).coordinate

            buf = io.BytesIO()
            wb.save(buf)
            return buf.getvalue()

        data = await asyncio.to_thread(_montar)
        nome = _nome_padrao(titulo, ecossistema, "xlsx", versao, self.valves.ECOSSISTEMA_PADRAO)
        ct = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        link = await _salvar_e_linkar(data, nome, ct, _get_user_id(__user__))
        return "Arquivo gerado com sucesso. " + link

    async def _docx(self, titulo, secoes, marca, __user__,
                    ecossistema="", versao=1, imagens=None):
        raw_secoes = secoes
        secoes = _itens_loose(secoes, _texto_para_secao)
        if not secoes:
            return _diag_entrada_vazia("gerar_docx", "secoes", raw_secoes)
        imgs = _normalizar_imagens(imagens)

        def _montar():
            from docx import Document
            from docx.shared import Pt, RGBColor, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            verde = RGBColor(*_hex_to_rgb(NIDUM_VERDE))
            terracota = RGBColor(*_hex_to_rgb(NIDUM_TERRACOTA))
            cinza = RGBColor(*_hex_to_rgb(NIDUM_CINZA))
            ink = RGBColor(*_hex_to_rgb(NIDUM_PRETO))

            doc = Document()

            # Logo de marca (regra 2.4.0: abertura E encerramento). Pagina branca ->
            # logo terracota. Centralizado, nunca em toda pagina - so topo e fim.
            def _logo_par(cor):
                lp = _logo_path(cor)
                if not lp:
                    return
                try:
                    par = doc.add_paragraph()
                    par.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    par.add_run().add_picture(lp, width=Inches(1.6))
                except Exception:
                    log.exception("gerador_nidum: falha ao inserir logo docx")

            if marca:
                try:
                    normal = doc.styles["Normal"]
                    normal.font.name = NIDUM_FONT
                    normal.font.size = Pt(11)
                    normal.font.color.rgb = ink
                    pf = normal.paragraph_format
                    pf.line_spacing = 1.32
                    pf.space_after = Pt(8)
                    for lvl, cor, sz in (
                        ("Title", verde, 30),
                        ("Heading 1", verde, 17),
                        ("Heading 2", terracota, 13),
                    ):
                        try:
                            st = doc.styles[lvl]
                            st.font.name = NIDUM_FONT_LOGO
                            st.font.color.rgb = cor
                            st.font.size = Pt(sz)
                        except Exception:
                            pass
                except Exception:
                    log.exception("gerador_nidum: falha ao aplicar estilos docx")

            def _imagem_par(img):
                # Imagem do usuario, centralizada, dentro da largura util da pagina.
                # add_picture com SO a largura ja preserva a proporcao (o python-docx
                # calcula a altura); depois conferimos a altura e, se estourar, reduzimos
                # os DOIS lados pelo mesmo fator - a forma nunca muda.
                LARG_MAX = Inches(5.9)    # A4 menos as margens padrao
                ALT_MAX = Inches(6.5)
                try:
                    par = doc.add_paragraph()
                    par.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    run = par.add_run()
                    pic = run.add_picture(io.BytesIO(img["bytes"]))
                    w, h = _encaixar(pic.width, pic.height, LARG_MAX, ALT_MAX)
                    pic.width = int(w)
                    pic.height = int(h)
                except Exception:
                    log.exception(
                        "gerador_nidum: falha ao inserir imagem do usuario no docx"
                    )

            if marca:
                _logo_par("terracota")   # logo de ABERTURA (topo)

            h = doc.add_heading(titulo or "Documento", level=0)
            for run in h.runs:
                run.font.name = NIDUM_FONT_LOGO
                if marca:
                    run.font.color.rgb = verde

            for sec in secoes:
                if sec.get("heading"):
                    hs = doc.add_heading(sec["heading"], level=1)
                    for run in hs.runs:
                        run.font.name = NIDUM_FONT_LOGO
                        if marca:
                            run.font.color.rgb = verde
                for par in _lista(sec.get("paragrafos")):
                    p = doc.add_paragraph(str(par))
                    for run in p.runs:
                        run.font.name = NIDUM_FONT
                        run.font.size = Pt(11)
                        run.font.color.rgb = ink
                for b in _lista(sec.get("bullets")):
                    pb = doc.add_paragraph(str(b), style="List Bullet")
                    for run in pb.runs:
                        run.font.name = NIDUM_FONT
                        run.font.color.rgb = ink
                img_s = _imagem_do_item(imgs, sec)
                if img_s is not None:
                    _imagem_par(img_s)

            if marca:
                _logo_par("terracota")   # logo de ENCERRAMENTO (fim)

            if marca:
                try:
                    fp = doc.sections[0].footer.paragraphs[0]
                    fr = fp.add_run("nidum. fazer da casa um ninho.")
                    fr.font.name = NIDUM_FONT
                    fr.font.size = Pt(8)
                    fr.font.color.rgb = ink
                    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                except Exception:
                    log.exception("gerador_nidum: falha ao aplicar rodape docx")

            buf = io.BytesIO()
            doc.save(buf)
            return buf.getvalue()

        data = await asyncio.to_thread(_montar)
        nome = _nome_padrao(titulo, ecossistema, "docx", versao, self.valves.ECOSSISTEMA_PADRAO)
        ct = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        link = await _salvar_e_linkar(data, nome, ct, _get_user_id(__user__))
        return "Arquivo gerado com sucesso. " + link

    async def _pdf(self, titulo, secoes, marca, __user__,
                   ecossistema="", versao=1, imagens=None):
        raw_secoes = secoes
        secoes = _itens_loose(secoes, _texto_para_secao)
        if not secoes:
            return _diag_entrada_vazia("gerar_pdf", "secoes", raw_secoes)
        imgs = _normalizar_imagens(imagens)

        def _montar():
            from reportlab.lib.pagesizes import A4
            from reportlab.lib import colors
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import mm
            from reportlab.platypus import (
                SimpleDocTemplate,
                Paragraph,
                Spacer,
                Table,
                TableStyle,
                ListFlowable,
                ListItem,
                Image,
                HRFlowable,
            )
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont

            verde = colors.HexColor("#" + NIDUM_VERDE)
            terracota = colors.HexColor("#" + NIDUM_TERRACOTA)
            preto = colors.HexColor("#" + NIDUM_PRETO)
            creme = colors.HexColor("#" + NIDUM_CREME)
            cremealt = colors.HexColor("#" + NIDUM_CREME_ALT)
            cinza = colors.HexColor("#" + NIDUM_CINZA)
            branco = colors.HexColor("#" + NIDUM_BRANCO)
            ink = preto
            linha_cor = colors.HexColor("#" + NIDUM_CINZA)

            # Registra a fonte da marca (fallback Helvetica se nao encontrar)
            FONT = "Helvetica"
            FONT_B = "Helvetica-Bold"
            FONT_TITULO = "Helvetica-Bold"
            try:
                reg = _font_path("MaximaNouva-Regular.ttf")
                bold = _font_path("MaximaNouva-Bold.ttf")
                xbold = _font_path("MaximaNouva-ExtraBold.ttf")
                ib = _font_path("Ibrand.ttf")
                if reg:
                    pdfmetrics.registerFont(TTFont("MaximaNouva", reg))
                    FONT = "MaximaNouva"
                if bold:
                    pdfmetrics.registerFont(TTFont("MaximaNouva-Bold", bold))
                    FONT_B = "MaximaNouva-Bold"
                # ExtraBold: o brandbook cita o peso e o arquivo existe no build desde
                # 13/07, mas nunca era registrado - ficava inerte no disco.
                if xbold:
                    pdfmetrics.registerFont(TTFont("MaximaNouva-ExtraBold", xbold))
                if reg and bold:
                    pdfmetrics.registerFontFamily(
                        "MaximaNouva", normal="MaximaNouva", bold="MaximaNouva-Bold"
                    )
                if ib:
                    pdfmetrics.registerFont(TTFont("Ibrand", ib))
                    FONT_TITULO = "Ibrand"
                else:
                    FONT_TITULO = FONT_B
            except Exception:
                log.exception("gerador_nidum: falha ao registrar fontes no pdf")
                FONT, FONT_B, FONT_TITULO = "Helvetica", "Helvetica-Bold", "Helvetica-Bold"

            styles = getSampleStyleSheet()
            st_titulo = ParagraphStyle(
                "NidumTitulo", parent=styles["Title"],
                textColor=verde if marca else ink, fontName=FONT_TITULO,
                fontSize=25, leading=29,
            )
            st_head = ParagraphStyle(
                "NidumHead", parent=styles["Heading1"],
                textColor=verde if marca else ink, fontName=FONT_TITULO,
                fontSize=15, leading=19,
            )
            st_body = ParagraphStyle(
                "NidumBody",
                parent=styles["BodyText"],
                textColor=ink,
                fontName=FONT,
                fontSize=11,
                leading=16,
            )

            buf = io.BytesIO()
            docp = SimpleDocTemplate(
                buf,
                pagesize=A4,
                leftMargin=22 * mm,
                rightMargin=22 * mm,
                topMargin=20 * mm,
                bottomMargin=18 * mm,
            )

            def _fundo(canvas, doc):
                if not marca:
                    return
                canvas.saveState()
                canvas.setFillColor(creme)
                canvas.rect(0, 0, doc.pagesize[0], doc.pagesize[1], fill=1, stroke=0)
                # rodape: assinatura + numero de pagina
                canvas.setFillColor(preto)
                canvas.setFont(FONT, 8)
                canvas.drawString(
                    22 * mm, 11 * mm, "nidum. fazer da casa um ninho."
                )
                canvas.drawRightString(
                    doc.pagesize[0] - 22 * mm, 11 * mm, str(doc.page)
                )
                canvas.restoreState()

            flow = []
            if marca:
                lp = _logo_path("terracota")
                if lp:
                    try:
                        flow.append(
                            Image(lp, width=40 * mm, height=22.5 * mm, hAlign="LEFT")
                        )
                        flow.append(Spacer(1, 2 * mm))
                    except Exception:
                        log.exception("gerador_nidum: falha ao inserir logo no pdf")
            flow.append(Paragraph(titulo or "Documento", st_titulo))
            flow.append(Spacer(1, 2 * mm))
            if marca:
                flow.append(
                    HRFlowable(
                        width="100%", thickness=1.5, color=terracota,
                        spaceBefore=1, spaceAfter=7 * mm, lineCap="round",
                    )
                )
            else:
                flow.append(Spacer(1, 6 * mm))

            for sec in secoes:
                if sec.get("heading"):
                    flow.append(Paragraph(sec["heading"], st_head))
                    flow.append(Spacer(1, 2 * mm))
                for par in _lista(sec.get("paragrafos")):
                    flow.append(Paragraph(str(par), st_body))
                    flow.append(Spacer(1, 2 * mm))
                bullets = _lista(sec.get("bullets"))
                if bullets:
                    items = [ListItem(Paragraph(str(b), st_body)) for b in bullets]
                    flow.append(ListFlowable(items, bulletType="bullet"))
                    flow.append(Spacer(1, 2 * mm))
                img_s = _imagem_do_item(imgs, sec)
                if img_s is not None:
                    # ImageReader da o tamanho NATIVO em pixels; _encaixar converte para
                    # pontos com fator unico. WEBP depende de PIL com suporte a webp - se
                    # a lib nao aceitar, cai no except, loga e o PDF sai SEM a imagem.
                    try:
                        from reportlab.lib.utils import ImageReader

                        ir = ImageReader(io.BytesIO(img_s["bytes"]))
                        nat_w, nat_h = ir.getSize()
                        larg_util = docp.width
                        w, h = _encaixar(nat_w, nat_h, larg_util, 150 * mm)
                        flow.append(Spacer(1, 3 * mm))
                        flow.append(
                            Image(io.BytesIO(img_s["bytes"]), width=w, height=h,
                                  hAlign="CENTER")
                        )
                        flow.append(Spacer(1, 4 * mm))
                    except Exception:
                        log.exception(
                            "gerador_nidum: falha ao inserir imagem do usuario no pdf"
                        )
                tabela = sec.get("tabela")
                if tabela:
                    tabela = [_lista(r) for r in _lista(tabela)]
                    t = Table(tabela, hAlign="LEFT")
                    estilo_t = [
                        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                        ("FONTSIZE", (0, 0), (-1, -1), 10),
                        ("TEXTCOLOR", (0, 1), (-1, -1), ink),
                        ("LEFTPADDING", (0, 0), (-1, -1), 9),
                        ("RIGHTPADDING", (0, 0), (-1, -1), 9),
                        ("TOPPADDING", (0, 0), (-1, -1), 7),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 7),
                        ("LINEBELOW", (0, 1), (-1, -1), 0.4, linha_cor),
                    ]
                    if marca:
                        estilo_t.append(("BACKGROUND", (0, 0), (-1, 0), verde))
                        estilo_t.append(("TEXTCOLOR", (0, 0), (-1, 0), creme))
                        estilo_t.append(("FONTNAME", (0, 0), (-1, 0), FONT_B))
                        estilo_t.append(("FONTNAME", (0, 1), (-1, -1), FONT))
                        estilo_t.append(("TOPPADDING", (0, 0), (-1, 0), 9))
                        estilo_t.append(("BOTTOMPADDING", (0, 0), (-1, 0), 9))
                        estilo_t.append(
                            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [branco, cremealt])
                        )
                    t.setStyle(TableStyle(estilo_t))
                    flow.append(t)
                    flow.append(Spacer(1, 4 * mm))

            # Logo de ENCERRAMENTO (regra 2.4.0: abertura E fim). O topo ja tem logo (acima);
            # aqui fecha o documento. Centralizado, menor, terracota (pagina areia).
            if marca:
                lp_fim = _logo_path("terracota")
                if lp_fim:
                    try:
                        flow.append(Spacer(1, 8 * mm))
                        flow.append(
                            Image(lp_fim, width=32 * mm, height=18 * mm, hAlign="CENTER")
                        )
                    except Exception:
                        log.exception("gerador_nidum: falha ao inserir logo de fim no pdf")

            docp.build(flow, onFirstPage=_fundo, onLaterPages=_fundo)
            return buf.getvalue()

        data = await asyncio.to_thread(_montar)
        nome = _nome_padrao(titulo, ecossistema, "pdf", versao, self.valves.ECOSSISTEMA_PADRAO)
        link = await _salvar_e_linkar(
            data, nome, "application/pdf", _get_user_id(__user__)
        )
        return "Arquivo gerado com sucesso. " + link
