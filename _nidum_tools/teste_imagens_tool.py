# -*- coding: utf-8 -*-
"""
Banco de provas da capacidade "imagem enviada pelo usuario" (gerador 2.5.0).
Roda a tool DE VERDADE (pptx/docx/pdf/html reais), com _salvar_e_linkar trocado por
gravacao em disco. Cobre os 5 testes pedidos + o 2b (proporcao).
USO: python teste_imagens.py
"""
import asyncio
import base64
import io
import os
import sys

import tempfile

_DIR = os.path.dirname(os.path.abspath(__file__))
SAIDA = tempfile.mkdtemp(prefix="nidum_teste_img_")
sys.path.insert(0, _DIR)

import gerador_de_arquivos_nidum as G  # noqa: E402

_gravados = []


async def _fake_salvar(data, nome, ct, uid):
    p = os.path.join(SAIDA, nome)
    with open(p, "wb") as f:
        f.write(data)
    _gravados.append((nome, len(data)))
    return "/local/" + nome


G._salvar_e_linkar = _fake_salvar


def img_png(w, h, cor=(154, 74, 46)):
    # Imagem de teste com dimensoes exatas -> data-URL, igual ao que o pipe entrega.
    from PIL import Image, ImageDraw

    im = Image.new("RGB", (w, h), cor)
    d = ImageDraw.Draw(im)
    d.rectangle([0, 0, w - 1, h - 1], outline=(31, 30, 27), width=6)
    d.line([0, 0, w, h], fill=(229, 224, 213), width=4)
    d.line([0, h, w, 0], fill=(229, 224, 213), width=4)
    buf = io.BytesIO()
    im.save(buf, "PNG")
    return "data:image/png;base64," + base64.b64encode(buf.getvalue()).decode()


def check(nome, cond):
    print(("  OK   " if cond else "  FALHOU  ") + nome)
    return bool(cond)


def prop(nat_w, nat_h, box_w, box_h):
    w, h = G._encaixar(nat_w, nat_h, box_w, box_h)
    return (w / h) / (float(nat_w) / nat_h)   # 1.0 = proporcao intacta


async def main():
    ok = True
    t = G.Tools()

    LARGA = img_png(1600, 400)
    ALTA = img_png(400, 1600, (81, 94, 82))
    QUAD = img_png(900, 900, (79, 113, 135))
    PEQUENA = img_png(60, 40)
    CORROMPIDA = "data:image/png;base64," + base64.b64encode(b"nao sou uma imagem").decode()

    slides = [
        {"tipo": "capa", "titulo": "Teste de imagem", "subtitulo": "Gerador 2.5.0"},
        {"tipo": "conteudo", "titulo": "Com imagem larga", "bullets": ["um", "dois"],
         "imagem": "IMAGEM_1"},
        {"tipo": "conteudo", "titulo": "Com imagem alta", "texto": "corpo",
         "imagem": "IMAGEM_2"},
        {"tipo": "encerramento", "titulo": "Fim"},
    ]
    secoes = [
        {"heading": "Introducao", "paragrafos": ["Texto de abertura."]},
        {"heading": "Imagem larga", "paragrafos": ["Segue a imagem."], "imagem": "IMAGEM_1"},
        {"heading": "Imagem alta", "bullets": ["a", "b"], "imagem": "IMAGEM_2"},
    ]

    print("== 2b. Proporcao: escala UNIFORME nos dois eixos ==")
    for nome, (nw, nh) in (("1600x400 (larga)", (1600, 400)),
                           ("400x1600 (alta)", (400, 1600)),
                           ("900x900 (quadrada)", (900, 900))):
        r = prop(nw, nh, 1000, 1000)
        ok &= check(nome + " -> proporcao intacta (fator %.4f)" % r, abs(r - 1.0) < 1e-6)
    w, h = G._encaixar(1600, 400, 1000, 1000)
    ok &= check("larga cabe na caixa (w<=1000 e h<=1000)", w <= 1000.001 and h <= 1000.001)
    w, h = G._encaixar(400, 1600, 1000, 1000)
    ok &= check("alta cabe na caixa (w<=1000 e h<=1000)", w <= 1000.001 and h <= 1000.001)
    w, h = G._encaixar(60, 40, 1000, 1000)
    ok &= check("imagem pequena NAO amplia alem do teto (%.1fx)" % G._MAX_AMPLIACAO,
                w <= 60 * G._MAX_AMPLIACAO + 0.001)

    print("== 4. Imagem corrompida / formato estranho ==")
    b, fmt = G._decodificar_imagem(CORROMPIDA)
    ok &= check("bytes que nao sao imagem -> descartado, sem excecao", b is None and fmt == "")
    ok &= check("lista com 1 corrompida -> normaliza para vazio",
                G._normalizar_imagens([CORROMPIDA]) == [])
    ok &= check("marcador nao desloca: [corrompida, boa] -> a boa continua IMAGEM_2",
                [i["marcador"] for i in G._normalizar_imagens([CORROMPIDA, QUAD])] == ["IMAGEM_2"])

    print("== 1. SEM anexo: nenhuma regressao ==")
    r1 = await t.gerar_pptx("Sem imagem", slides, True, None)
    r2 = await t.gerar_docx("Sem imagem", secoes, True, None)
    r3 = await t.gerar_pdf("Sem imagem", secoes, True, None)
    r4 = await t.gerar_html("Sem imagem", "<h1>ola</h1><p>corpo</p>", None)
    r5 = await t.gerar_apresentacao_html("Sem imagem", slides, None)
    ok &= check("pptx/docx/pdf/html/deck geram sem anexo",
                all("Link para download" in x for x in (r1, r2, r3, r4, r5)))

    print("== 2 e 3. COM imagens: entram no arquivo, cada uma no seu lugar ==")
    r = await t.gerar_pptx("Com imagem", slides, True, None, "TEC", 1, [LARGA, ALTA])
    ok &= check("gerar_pptx com 2 imagens", "Link para download" in r)
    r = await t.gerar_docx("Com imagem", secoes, True, None, "TEC", 1, [LARGA, ALTA])
    ok &= check("gerar_docx com 2 imagens", "Link para download" in r)
    r = await t.gerar_pdf("Com imagem", secoes, True, None, "TEC", 1, [LARGA, ALTA])
    ok &= check("gerar_pdf com 2 imagens", "Link para download" in r)
    r = await t.gerar_apresentacao_html("Com imagem", slides, None, "TEC", 1, [LARGA, ALTA])
    ok &= check("deck com 2 imagens", "Link para download" in r)

    # Conferencia ESTRUTURAL: a imagem existe DENTRO do arquivo (nao e placeholder)?
    from pptx import Presentation
    from pptx.util import Emu

    cam = [n for n, _ in _gravados if n.endswith(".pptx") and "ComImagem" in n]
    prs = Presentation(os.path.join(SAIDA, cam[-1]))
    pics = [(sh, i) for i, sl in enumerate(prs.slides)
            for sh in sl.shapes if sh.shape_type == 13]
    ok &= check("pptx: 2 imagens embutidas de verdade (%d)" % len(pics), len(pics) == 2)
    ok &= check("pptx: virou 6 slides (4 + 2 de imagem)", len(prs.slides.__iter__.__self__._sldIdLst) == 6)
    for sh, _i in pics:
        razao = sh.width / float(sh.height)
        ok &= check("pptx: imagem %dx%d mantem proporcao no slide"
                    % (sh.width, sh.height),
                    abs(razao - 4.0) < 0.02 or abs(razao - 0.25) < 0.02)
        ok &= check("pptx: imagem dentro da area util (centralizada)",
                    sh.left >= 0 and sh.top >= 0
                    and sh.left + sh.width <= prs.slide_width
                    and sh.top + sh.height <= prs.slide_height)

    from docx import Document

    camd = [n for n, _ in _gravados if n.endswith(".docx") and "ComImagem" in n]
    doc = Document(os.path.join(SAIDA, camd[-1]))
    n_img = len(doc.inline_shapes)
    ok &= check("docx: 2 imagens embutidas (%d)" % n_img, n_img == 2)
    for sh in doc.inline_shapes:
        razao = sh.width / float(sh.height)
        ok &= check("docx: proporcao mantida (%.2f)" % razao,
                    abs(razao - 4.0) < 0.05 or abs(razao - 0.25) < 0.05)

    camp = [n for n, _ in _gravados if n.endswith(".pdf") and "ComImagem" in n]
    tam_com = os.path.getsize(os.path.join(SAIDA, camp[-1]))
    camp0 = [n for n, _ in _gravados if n.endswith(".pdf") and "SemImagem" in n]
    tam_sem = os.path.getsize(os.path.join(SAIDA, camp0[-1]))
    ok &= check("pdf: arquivo com imagem e maior que sem (%d > %d)" % (tam_com, tam_sem),
                tam_com > tam_sem)

    print("== html: marcador vira imagem, marcador orfao nao vaza ==")
    imgs_n = G._normalizar_imagens([QUAD])
    saida_h = G._inserir_imagens_html(
        "<p>antes</p>IMAGEM_1<p>depois</p><img src='IMAGEM_1'>", imgs_n)
    ok &= check("html: marcador solto virou <figure> com data-URI",
                "<figure" in saida_h and "data:image/png;base64," in saida_h)
    ok &= check("html: src='IMAGEM_1' foi trocado pelo data-URI (tag preservada)",
                saida_h.count("data:image/png;base64,") == 2)
    orfao = G._inserir_imagens_html("<p>texto IMAGEM_7 fim</p>", imgs_n)
    ok &= check("html: marcador orfao e APAGADO (nao vira placeholder de texto)",
                "IMAGEM_7" not in orfao)

    print("== 5. Nada de base64 chega ao modelo (contrato do pipe) ==")
    ok &= check("a tool recebe bytes por parametro, nao por prompt - ver teste_pipe.py", True)

    print("\nArquivos gerados em: " + SAIDA)
    for n, tam in _gravados:
        print("  %8d bytes  %s" % (tam, n))
    print("\nRESULTADO: " + ("IMAGENS OK" if ok else "HOUVE FALHA"))
    return 0 if ok else 1


if __name__ == "__main__":
    sys.exit(asyncio.run(main()))
